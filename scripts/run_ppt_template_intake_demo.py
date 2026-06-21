from __future__ import annotations

import json
from pathlib import Path
import shutil
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt

from drawai.ppt_template_intake import (
    apply_minimal_fill_plan,
    build_prisma_reference_style_spec,
    build_slot_schema_preview,
    build_template_reference_payload,
    create_minimal_fill_plan,
    save_json,
    write_ppt_template_spec,
)


REPO_ROOT = Path(__file__).resolve().parents[1]
INTAKE_DIR = REPO_ROOT / "outputs" / "ppt_template_intake_demo"
SHOWCASE_DIR = REPO_ROOT / "outputs" / "ppt_master_alignment_showcase"
REFERENCE_TEST_DIR = REPO_ROOT / "outputs" / "codex_reference_image_generation_test"
PRISMA_TEMPLATE_MANIFEST = REPO_ROOT / "templates" / "slide_image" / "prisma_flow_diagram" / "template.json"
ORIGINAL_PRISMA_IMAGE = Path(r"C:\Users\yanrupeng\AppData\Local\hermes\image_cache\img_3801f9a210be.jpg")


def main() -> None:
    INTAKE_DIR.mkdir(parents=True, exist_ok=True)
    SHOWCASE_DIR.mkdir(parents=True, exist_ok=True)

    demo_pptx = INTAKE_DIR / "demo_template.pptx"
    _build_demo_template_pptx(demo_pptx)

    template_spec = write_ppt_template_spec(demo_pptx, INTAKE_DIR / "template_spec.json")
    slot_preview = build_slot_schema_preview(template_spec)
    save_json(slot_preview, INTAKE_DIR / "slot_schema_preview.json")

    reference_source = _resolve_prisma_source()
    reference_style_spec = build_prisma_reference_style_spec(
        source_image_path=reference_source,
        original_source_image_path=ORIGINAL_PRISMA_IMAGE,
        template_manifest_path=PRISMA_TEMPLATE_MANIFEST,
    )
    save_json(reference_style_spec, REFERENCE_TEST_DIR / "reference_style_spec.json")

    edit_payload = build_template_reference_payload(
        template_spec=template_spec,
        reference_style_spec=reference_style_spec,
        user_topic="DrawAI PPT 图像生成能力验证流程",
        language="zh",
    )
    fill_plan = create_minimal_fill_plan(template_spec, user_topic="DrawAI PPT 模板 intake 与参考图 spec 对齐")
    save_json(fill_plan, INTAKE_DIR / "fill_plan.json")
    fill_result = apply_minimal_fill_plan(demo_pptx, fill_plan, INTAKE_DIR / "output_demo.pptx")
    save_json(fill_result, INTAKE_DIR / "fill_result.json")

    showcase_files = _copy_showcase_files(
        demo_pptx=demo_pptx,
        template_spec=template_spec,
        slot_preview=slot_preview,
        reference_style_spec=reference_style_spec,
        edit_payload=edit_payload,
        fill_plan=fill_plan,
        fill_result=fill_result,
    )
    summary = _build_summary(showcase_files=showcase_files, template_spec=template_spec, fill_result=fill_result)
    save_json(summary, SHOWCASE_DIR / "summary.json")
    (SHOWCASE_DIR / "summary.md").write_text(_summary_md(summary), encoding="utf-8")
    _write_contact_sheet(summary, SHOWCASE_DIR / "contact_sheet.jpg")

    print(json.dumps(summary, ensure_ascii=False, indent=2))


def _build_demo_template_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    _add_cover_slide(prs, blank)
    _add_content_slide(prs, blank)
    _add_flow_slide(prs, blank)
    _add_data_slide(prs, blank)

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))


def _add_cover_slide(prs: Presentation, blank: Any) -> None:
    slide = prs.slides.add_slide(blank)
    _add_header_bar(slide, "模板 PPTX 输入", "design_spec / spec_lock")
    _add_textbox(slide, "slot_cover_title", "真实模板 intake", 0.75, 1.25, 7.5, 0.75, 34, True, "111827")
    _add_textbox(
        slide,
        "slot_cover_subtitle",
        "从示例 deck 抽取 slide size、theme、slot schema，再进行最小可编辑填槽",
        0.78,
        2.1,
        7.6,
        0.6,
        17,
        False,
        "374151",
    )
    for idx, (label, value) in enumerate([("Slide", "4"), ("Slots", "auto"), ("Output", "PPTX")]):
        left = 0.8 + idx * 2.1
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(4.85), Inches(1.72), Inches(0.92))
        card.name = f"slot_metric_{idx + 1}"
        _style_shape(card, fill="F8FAFC", line="CBD5E1")
        _set_text(card, f"{label}\n{value}", 16, True, "111827")
    _add_textbox(slide, "slot_cover_note", "Native PPTX, not a flat image", 9.35, 5.85, 2.7, 0.35, 13, False, "64748B")


def _add_content_slide(prs: Presentation, blank: Any) -> None:
    slide = prs.slides.add_slide(blank)
    _add_header_bar(slide, "内容结构", "slot fill")
    _add_textbox(slide, "slot_content_title", "PPT-master 对齐的数据流", 0.65, 0.78, 6.6, 0.55, 25, True, "111827")
    cards = [
        ("slot_card_1", "Template Deck", "读取用户 PPTX / 示例 deck，保留原生布局与可编辑元素。"),
        ("slot_card_2", "Template Spec", "抽取 slide_size、theme、slots、tables、charts 与 role_guess。"),
        ("slot_card_3", "Fill Plan", "用稳定 slot_id 替换内容，避免把模板降级成 prompt。"),
    ]
    for idx, (name, title, body) in enumerate(cards):
        left = 0.75 + idx * 4.1
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.0), Inches(3.45), Inches(2.55))
        box.name = name
        _style_shape(box, fill="FFFFFF", line="CBD5E1")
        _set_text(box, f"{title}\n{body}", 16, True if idx == 0 else False, "111827")
    _add_textbox(slide, "slot_content_caption", "第一阶段：先证明 intake + fill 可跑通；后续扩展 Slide Master 与真实模板库。", 0.8, 5.55, 11.2, 0.5, 14, False, "475569")


def _add_flow_slide(prs: Presentation, blank: Any) -> None:
    slide = prs.slides.add_slide(blank)
    _add_header_bar(slide, "流程页", "layout slots")
    _add_textbox(slide, "slot_flow_title", "参考图 / PPT 模板 / 数据源如何合并", 0.65, 0.75, 7.2, 0.5, 24, True, "111827")
    steps = ["模板 PPTX", "Style Spec", "Slot Fill", "DrawAI 重建"]
    for idx, label in enumerate(steps):
        left = 0.9 + idx * 3.05
        node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(3.0), Inches(2.25), Inches(0.9))
        node.name = f"slot_flow_step_{idx + 1}"
        _style_shape(node, fill="FFFFFF", line="111827")
        _set_text(node, label, 16, True, "111827")
        if idx < len(steps) - 1:
            arrow = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(left + 2.28), Inches(3.45), Inches(left + 2.9), Inches(3.45))
            arrow.line.color.rgb = RGBColor.from_string("111827")
            arrow.line.width = Pt(1.4)
    _add_textbox(slide, "slot_flow_note", "spec_lock 控制颜色/字体/槽位；fill_plan 只替换允许变动的内容。", 1.0, 5.3, 9.6, 0.5, 15, False, "475569")


def _add_data_slide(prs: Presentation, blank: Any) -> None:
    slide = prs.slides.add_slide(blank)
    _add_header_bar(slide, "验证页", "tables / charts")
    _add_textbox(slide, "slot_data_title", "第一阶段可验证能力", 0.65, 0.76, 5.8, 0.5, 24, True, "111827")
    table_shape = slide.shapes.add_table(4, 3, Inches(0.8), Inches(1.75), Inches(5.7), Inches(2.25))
    table_shape.name = "native_table_capability_matrix"
    table = table_shape.table
    rows = [
        ["能力", "状态", "证据"],
        ["PPTX intake", "已跑通", "template_spec.json"],
        ["参考图 spec", "已结构化", "reference_style_spec.json"],
        ["填槽 PPTX", "最小 demo", "output_demo.pptx"],
    ]
    for row_index, row in enumerate(rows):
        for col_index, text in enumerate(row):
            table.cell(row_index, col_index).text = text
    chart_data = ChartData()
    chart_data.categories = ["intake", "spec", "fill"]
    chart_data.add_series("coverage", (90, 70, 45))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(7.1),
        Inches(1.65),
        Inches(5.05),
        Inches(2.7),
        chart_data,
    )
    chart.name = "native_chart_phase_coverage"
    _add_textbox(slide, "slot_data_caption", "图表值是 demo 覆盖度，不代表产品完成度；用于验证 chart inventory。", 0.85, 5.55, 11.0, 0.45, 13, False, "64748B")


def _add_header_bar(slide: Any, left_text: str, right_text: str) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.32))
    bar.name = "slot_header_bar"
    _style_shape(bar, fill="FBBF24", line="FBBF24")
    _set_text(bar, f"{left_text}    {right_text}", 11, True, "111827")


def _add_textbox(
    slide: Any,
    name: str,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int,
    bold: bool,
    color: str,
) -> Any:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    shape.name = name
    _set_text(shape, text, size, bold, color)
    return shape


def _set_text(shape: Any, text: str, size: int, bold: bool, color: str) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    text_frame.word_wrap = True


def _style_shape(shape: Any, *, fill: str, line: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.color.rgb = RGBColor.from_string(line)
    shape.line.width = Pt(1)


def _resolve_prisma_source() -> Path:
    reference_copy = REFERENCE_TEST_DIR / "source_reference.jpg"
    if reference_copy.is_file():
        return reference_copy
    return ORIGINAL_PRISMA_IMAGE


def _copy_showcase_files(
    *,
    demo_pptx: Path,
    template_spec: dict[str, Any],
    slot_preview: dict[str, Any],
    reference_style_spec: dict[str, Any],
    edit_payload: dict[str, Any],
    fill_plan: dict[str, Any],
    fill_result: dict[str, Any],
) -> dict[str, str]:
    files = {
        "demo_template_pptx": SHOWCASE_DIR / "demo_template.pptx",
        "template_spec_json": SHOWCASE_DIR / "template_spec.json",
        "slot_schema_preview_json": SHOWCASE_DIR / "slot_schema_preview.json",
        "reference_style_spec_json": SHOWCASE_DIR / "reference_style_spec.json",
        "template_reference_edit_payload_json": SHOWCASE_DIR / "template_reference_edit_payload.json",
        "fill_plan_json": SHOWCASE_DIR / "fill_plan.json",
        "fill_result_json": SHOWCASE_DIR / "fill_result.json",
        "output_demo_pptx": SHOWCASE_DIR / "output_demo.pptx",
    }
    shutil.copy2(demo_pptx, files["demo_template_pptx"])
    shutil.copy2(INTAKE_DIR / "output_demo.pptx", files["output_demo_pptx"])
    save_json(template_spec, files["template_spec_json"])
    save_json(slot_preview, files["slot_schema_preview_json"])
    save_json(reference_style_spec, files["reference_style_spec_json"])
    save_json(edit_payload, files["template_reference_edit_payload_json"])
    save_json(fill_plan, files["fill_plan_json"])
    save_json(fill_result, files["fill_result_json"])
    return {key: str(value.resolve(strict=False)) for key, value in files.items()}


def _build_summary(*, showcase_files: dict[str, str], template_spec: dict[str, Any], fill_result: dict[str, Any]) -> dict[str, Any]:
    return {
        "schema": "drawai.ppt_master_alignment_showcase.summary.v1",
        "status": "ok",
        "output_dir": str(SHOWCASE_DIR.resolve(strict=False)),
        "intake_output_dir": str(INTAKE_DIR.resolve(strict=False)),
        "cases": [
            {
                "id": "A",
                "name": "demo PPTX -> template_spec.json",
                "status": "ok",
                "output": showcase_files["template_spec_json"],
                "observed": {
                    "slide_count": template_spec.get("slide_count"),
                    "slide_size": template_spec.get("slide_size"),
                    "layout_roles": template_spec.get("design_tokens", {}).get("layout_roles", {}),
                },
            },
            {
                "id": "B",
                "name": "template_spec -> slot schema / role_guess",
                "status": "ok",
                "output": showcase_files["slot_schema_preview_json"],
            },
            {
                "id": "C",
                "name": "PRISMA reference image -> reference_style_spec",
                "status": "ok",
                "output": showcase_files["reference_style_spec_json"],
            },
            {
                "id": "D",
                "name": "template_spec + reference_style_spec + user topic -> edit payload",
                "status": "ok",
                "output": showcase_files["template_reference_edit_payload_json"],
            },
            {
                "id": "E",
                "name": "minimal template-fill -> editable output_demo.pptx",
                "status": "ok" if not fill_result.get("missing_targets") else "warning",
                "output": showcase_files["output_demo_pptx"],
                "fill_result": fill_result,
            },
        ],
        "files": showcase_files,
        "limitations": [
            "This is a first-phase PPT-master alignment, not full PPT-master parity.",
            "Slide Master parsing is not complete; current spec is slide-bound layout/slot inventory.",
            "The minimal fill demo replaces text in existing slides; it does not yet clone/reorder template pages or edit native chart data.",
            "Reference image style spec is manually declared from the verified PRISMA reference; automatic visual token extraction is a next phase.",
        ],
    }


def _summary_md(summary: dict[str, Any]) -> str:
    lines = [
        "# PPT-master alignment showcase",
        "",
        "本轮展示从 prompt-only 纠偏到真实 PPTX/template intake 与最小可编辑填槽。",
        "",
        f"- 输出目录: `{summary['output_dir']}`",
        f"- Intake 输出目录: `{summary['intake_output_dir']}`",
        "",
        "## Cases",
    ]
    for case in summary["cases"]:
        lines.append(f"- {case['id']}. {case['name']}: {case['status']} -> `{case['output']}`")
    lines.extend(
        [
            "",
            "## 当前边界",
            "- 已能读取 demo PPTX 并输出 slide_size、theme、layout、slot、table、chart 粗略信息。",
            "- 已把 PRISMA 参考图表达为 layout/style/color/typography/content reference roles。",
            "- 已生成可编辑 `output_demo.pptx`，证明 slot fill 雏形可落地。",
            "- 尚未完成完整 Slide Master 解析、模板页克隆/重排、图表数据写回和视觉溢出自动检查。",
        ]
    )
    return "\n".join(lines) + "\n"


def _write_contact_sheet(summary: dict[str, Any], path: Path) -> None:
    width, height = 1800, 1080
    image = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(image)
    title_font = _font(42)
    body_font = _font(26)
    small_font = _font(22)
    draw.rectangle((0, 0, width, 120), fill=(251, 191, 36))
    draw.text((48, 34), "DrawAI PPT-master 对齐功能展示", fill=(17, 24, 39), font=title_font)
    y = 165
    for case in summary["cases"]:
        draw.rounded_rectangle((48, y, width - 48, y + 120), radius=16, outline=(203, 213, 225), width=2, fill=(248, 250, 252))
        draw.text((82, y + 24), f"{case['id']}. {case['name']}", fill=(15, 23, 42), font=body_font)
        draw.text((82, y + 68), f"status={case['status']}  output={Path(case['output']).name}", fill=(71, 85, 105), font=small_font)
        y += 145
    draw.text((48, height - 92), "注：该 contact sheet 是功能报告图，不是 PPT slide 渲染缩略图；可编辑 PPTX 见 output_demo.pptx。", fill=(100, 116, 139), font=small_font)
    path.parent.mkdir(parents=True, exist_ok=True)
    image.save(path, quality=92)


def _font(size: int) -> ImageFont.ImageFont:
    for font_path in [
        Path(r"C:\Windows\Fonts\msyh.ttc"),
        Path(r"C:\Windows\Fonts\simhei.ttf"),
        Path(r"C:\Windows\Fonts\arial.ttf"),
    ]:
        if font_path.is_file():
            return ImageFont.truetype(str(font_path), size)
    return ImageFont.load_default()


if __name__ == "__main__":
    main()
