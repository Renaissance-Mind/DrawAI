from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from drawai.ppt_template_intake import (
    PPT_TEMPLATE_SPEC_SCHEMA,
    REFERENCE_STYLE_SPEC_SCHEMA,
    apply_minimal_fill_plan,
    build_prisma_reference_style_spec,
    build_slot_schema_preview,
    create_minimal_fill_plan,
    intake_ppt_template,
    write_ppt_template_spec,
)


def test_ppt_template_intake_generates_template_spec(tmp_path: Path) -> None:
    source = tmp_path / "template.pptx"
    _make_demo_pptx(source)

    spec = intake_ppt_template(source)

    assert spec["schema"] == PPT_TEMPLATE_SPEC_SCHEMA
    assert spec["slide_count"] == 2
    assert spec["slide_size"]["width_emu"] > 0
    assert spec["layouts"][0]["slots"]
    assert spec["layouts"][0]["slots"][0]["slot_id"].startswith("s01_shape_")
    assert spec["spec_lock"]["lock_slot_geometry"] is True


def test_ppt_template_intake_writes_spec_and_slot_preview(tmp_path: Path) -> None:
    source = tmp_path / "template.pptx"
    output = tmp_path / "template_spec.json"
    _make_demo_pptx(source)

    spec = write_ppt_template_spec(source, output)
    preview = build_slot_schema_preview(spec)

    assert output.is_file()
    assert preview["layouts"]
    assert preview["layouts"][0]["slots"]
    assert preview["layouts"][0]["slots"][0]["role"] in {"title", "body", "label", "caption", "subtitle"}


def test_reference_style_spec_has_typed_reference_roles(tmp_path: Path) -> None:
    source_image = tmp_path / "reference.jpg"
    source_image.write_bytes(b"fake image placeholder")

    spec = build_prisma_reference_style_spec(source_image_path=source_image, original_source_image_path=source_image)

    assert spec["schema"] == REFERENCE_STYLE_SPEC_SCHEMA
    roles = {item["role"] for item in spec["reference_roles"]}
    assert {"layout_reference", "style_reference", "color_reference", "typography_reference", "content_reference"}.issubset(roles)
    assert spec["spec_lock"]["lock_layout_archetype"] is True
    assert "flow_nodes" in spec["design_tokens"]["slot_schema"]


def test_minimal_fill_plan_outputs_editable_pptx(tmp_path: Path) -> None:
    source = tmp_path / "template.pptx"
    output = tmp_path / "output_demo.pptx"
    _make_demo_pptx(source)
    spec = intake_ppt_template(source)
    fill_plan = create_minimal_fill_plan(spec, user_topic="DrawAI 模板填槽测试")

    result = apply_minimal_fill_plan(source, fill_plan, output)

    assert output.is_file()
    assert result["replaced_slot_count"] >= 2
    reopened = Presentation(str(output))
    slide_text = "\n".join(shape.text for slide in reopened.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False))
    assert "DrawAI" in slide_text


def _make_demo_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    _add_textbox(slide, "slot_title", "Demo Title", 0.8, 0.7, 5.4, 0.7, 30, True)
    _add_textbox(slide, "slot_subtitle", "Demo subtitle", 0.85, 1.55, 6.2, 0.45, 16, False)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.6), Inches(3.2), Inches(1.2))
    card.name = "slot_card"
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    card.line.color.rgb = RGBColor.from_string("CBD5E1")
    card.text = "Body placeholder"

    slide = prs.slides.add_slide(blank)
    _add_textbox(slide, "slot_data_title", "Data Title", 0.8, 0.7, 5.4, 0.7, 28, True)
    table_shape = slide.shapes.add_table(2, 2, Inches(0.8), Inches(1.8), Inches(4.0), Inches(1.1))
    table_shape.name = "native_table"
    table_shape.table.cell(0, 0).text = "Metric"
    table_shape.table.cell(0, 1).text = "Value"

    prs.save(str(path))


def _add_textbox(
    slide: object,
    name: str,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int,
    bold: bool,
) -> None:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    shape.name = name
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Microsoft YaHei"
    run.font.color.rgb = RGBColor.from_string("111827")
