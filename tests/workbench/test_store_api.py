from __future__ import annotations

import base64
import json
import io
import shutil
import threading
import time
import zipfile
from pathlib import Path

import pytest
import yaml
from fastapi.testclient import TestClient
from PIL import Image
from pptx import Presentation

from drawai.codex_python_sdk_imagegen import CodexGeneratedImage, CodexImageGenResult
from drawai.http_utils import is_loopback_url
from drawai.rmbg_client import RmbgResult
from drawai.v2.packages import write_asset_package, write_element_plan
from drawai.v2.schema import AssetPackage, ElementPlan, ProcessingIntent, RUN_PACKAGE_SCHEMA
import drawai.workbench.api as workbench_api
from drawai.workbench.api import create_app
from drawai.workbench.models import WorkbenchSettings
from drawai.workbench.runner import WorkbenchRunner, _archive_fs_path, _pipeline_failure_message, create_case_config
from drawai.workbench.store import WorkbenchStore


def test_store_creates_batch_case_stage_and_artifact(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    config_path = tmp_path / "config.yaml"
    source_path = tmp_path / "source.png"
    config_path.write_text("input: {}\n", encoding="utf-8")
    source_path.write_bytes(b"png")

    batch = store.create_batch(
        name="demo",
        input_mode="local_dir",
        max_concurrent_cases=2,
        auto_run_svg_after_analysis=True,
        config_path=config_path,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="case_001",
        source_image_path=source_path,
        config_path=config_path,
    )
    assert store.case_counts(batch.batch_id) == {"queued": 1}

    store.update_case_status(case.case_id, status="analysis_running", phase="analysis", stage="prepare")
    updated = store.get_case(case.case_id)
    assert updated.status == "analysis_running"
    assert updated.stage == "prepare"

    stage = store.start_stage_run(case.case_id, "prepare")
    store.finish_stage_run(stage.stage_run_id, status="ok")
    assert store.list_stage_runs(case.case_id)[0].status == "ok"

    artifact_path = Path(case.run_root) / "reports" / "status.json"
    artifact_path.parent.mkdir(parents=True)
    artifact_path.write_text("{}", encoding="utf-8")
    artifact = store.register_artifact(case.case_id, label="status", path=artifact_path, media_type="application/json")
    resolved = store.resolve_artifact(artifact.artifact_token)
    assert resolved.path == str(artifact_path)
    assert resolved.to_api()["url"].startswith("/api/artifacts/")


def test_pipeline_failure_message_prefers_exception_message() -> None:
    summary = {
        "status": "failed",
        "exception": {"message": "detailed pipeline failure"},
        "error": {"message": "legacy error"},
    }

    assert _pipeline_failure_message(summary) == "detailed pipeline failure"


@pytest.mark.parametrize(
    "existing_analysis",
    [
        None,
        {
            "schema": "drawai.codex_element_analysis.v1",
            "source": "v2.refined_elements",
            "elements": [],
        },
    ],
)
def test_runner_refine_stage_prepares_codex_analysis_when_required(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    existing_analysis: dict[str, object] | None,
) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )
    case_root = Path(case.run_root)
    config_path = create_case_config(
        base_config_path=base_config,
        source_image=source,
        output_dir=case_root,
        target_path=case_root / "drawai.config.yaml",
    )
    store.update_case_config_path(case.case_id, config_path)
    case = store.get_case(case.case_id)
    if existing_analysis is not None:
        _write_json(
            Path(case.run_root) / "reports" / "element_analysis_codex" / "element_analysis.json",
            existing_analysis,
        )
    codex_analysis_roots: list[Path] = []

    def fake_codex_analysis(cfg, paths) -> None:
        codex_analysis_roots.append(paths.root)
        _write_json(
            paths.element_analysis_json,
            {
                "schema": "drawai.codex_element_analysis.v1",
                "elements": [
                    {
                        "box_id": "B001",
                        "source_candidate_ids": ["B001"],
                        "bbox": [1, 1, 10, 10],
                        "category": "crop",
                        "type": "image",
                    }
                ],
            },
        )

    def fake_run_from_stage(
        config_path_or_config,
        from_stage: str,
        *,
        to_stage: str | None = None,
        **kwargs,
    ) -> dict[str, object]:
        assert config_path_or_config == case.config_path
        assert from_stage == "refine_elements"
        assert to_stage == "refine_elements"
        assert codex_analysis_roots == [Path(case.run_root)]
        return {"status": "ok", "artifacts": {}}

    monkeypatch.setattr("drawai.pipeline._run_codex_run0_asset_analysis", fake_codex_analysis)
    monkeypatch.setattr("drawai.workbench.runner.run_drawai_pipeline_from_stage", fake_run_from_stage)
    runner = WorkbenchRunner(store, _settings(tmp_path, base_config))

    runner._run_stage(case.case_id, "refine_elements")

    assert codex_analysis_roots == [Path(case.run_root)]
    saved_analysis_path = Path(case.run_root) / "reports" / "element_analysis_codex" / "element_analysis.json"
    saved_analysis = json.loads(saved_analysis_path.read_text(encoding="utf-8"))
    assert "source" not in saved_analysis


def test_runner_refine_stage_reuses_existing_external_codex_analysis(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )
    case_root = Path(case.run_root)
    config_path = create_case_config(
        base_config_path=base_config,
        source_image=source,
        output_dir=case_root,
        target_path=case_root / "drawai.config.yaml",
    )
    store.update_case_config_path(case.case_id, config_path)
    case = store.get_case(case.case_id)
    _write_json(
        Path(case.run_root) / "reports" / "element_analysis_codex" / "element_analysis.json",
        {"schema": "drawai.codex_element_analysis.v1", "elements": []},
    )

    def unexpected_codex_analysis(cfg, paths) -> None:
        raise AssertionError("existing external Codex analysis should be reused")

    def fake_run_from_stage(
        config_path_or_config,
        from_stage: str,
        *,
        to_stage: str | None = None,
        **kwargs,
    ) -> dict[str, object]:
        assert config_path_or_config == case.config_path
        assert from_stage == "refine_elements"
        assert to_stage == "refine_elements"
        return {"status": "ok", "artifacts": {}}

    monkeypatch.setattr("drawai.pipeline._run_codex_run0_asset_analysis", unexpected_codex_analysis)
    monkeypatch.setattr("drawai.workbench.runner.run_drawai_pipeline_from_stage", fake_run_from_stage)
    runner = WorkbenchRunner(store, _settings(tmp_path, base_config))

    runner._run_stage(case.case_id, "refine_elements")


def test_runner_refine_stage_regenerates_incomplete_codex_analysis(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )
    case_root = Path(case.run_root)
    config_path = create_case_config(
        base_config_path=base_config,
        source_image=source,
        output_dir=case_root,
        target_path=case_root / "drawai.config.yaml",
    )
    store.update_case_config_path(case.case_id, config_path)
    case = store.get_case(case.case_id)
    analysis_path = Path(case.run_root) / "reports" / "element_analysis_codex" / "element_analysis.json"
    _write_json(analysis_path, {"schema": "drawai.codex_element_analysis.v1", "elements": []})
    _write_json(analysis_path.parent / "run_status.json", {"status": "running"})
    codex_analysis_roots: list[Path] = []

    def fake_codex_analysis(cfg, paths) -> None:
        codex_analysis_roots.append(paths.root)
        _write_json(
            paths.element_analysis_json,
            {"schema": "drawai.codex_element_analysis.v1", "elements": []},
        )
        _write_json(paths.element_analysis_json.parent / "run_status.json", {"status": "ok"})

    def fake_run_from_stage(
        config_path_or_config,
        from_stage: str,
        *,
        to_stage: str | None = None,
        **kwargs,
    ) -> dict[str, object]:
        assert config_path_or_config == case.config_path
        assert from_stage == "refine_elements"
        assert to_stage == "refine_elements"
        assert codex_analysis_roots == [Path(case.run_root)]
        return {"status": "ok", "artifacts": {}}

    monkeypatch.setattr("drawai.pipeline._run_codex_run0_asset_analysis", fake_codex_analysis)
    monkeypatch.setattr("drawai.workbench.runner.run_drawai_pipeline_from_stage", fake_run_from_stage)
    runner = WorkbenchRunner(store, _settings(tmp_path, base_config))

    runner._run_stage(case.case_id, "refine_elements")

    assert codex_analysis_roots == [Path(case.run_root)]


def test_store_rejects_artifacts_outside_case_root(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    config_path = tmp_path / "config.yaml"
    source_path = tmp_path / "source.png"
    config_path.write_text("input: {}\n", encoding="utf-8")
    source_path.write_bytes(b"png")
    batch = store.create_batch(
        name="demo",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=config_path,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="case_001",
        source_image_path=source_path,
        config_path=config_path,
    )

    with pytest.raises(ValueError, match="outside case root"):
        store.register_artifact(case.case_id, label="secret", path=tmp_path / "secret.txt")

    with pytest.raises(ValueError, match="outside case root"):
        store.write_case_json(case.case_id, "../escape.json", {"bad": True})


def test_runner_completes_analysis_and_stops_for_asset_review(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="local_dir",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case_root = store.runs_root / batch.batch_id / "case_seed"
    config_path = create_case_config(
        base_config_path=base_config,
        source_image=source,
        output_dir=case_root,
        target_path=case_root / "drawai.config.yaml",
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=config_path,
    )
    observed_stages: list[str] = []

    def recording_stage_executor(case, stage: str) -> None:
        observed_stages.append(stage)
        _deterministic_stage_executor(case, stage)

    runner = WorkbenchRunner(store, _settings(tmp_path, base_config), stage_executor=recording_stage_executor)

    runner.submit_batch(batch.batch_id)
    runner.wait_for_idle(timeout=5)

    updated = store.get_case(case.case_id)
    assert updated.status == "assets_review"
    assert updated.stage == "plan_assets"
    assert observed_stages == [
        "prepare",
        "parse_elements",
        "fuse_elements",
        "refine_elements",
        "plan_assets",
    ]
    assert (Path(updated.run_root) / "reports" / "workbench" / "asset_draft.json").exists()
    assert (Path(updated.run_root) / "drawai_package.json").exists()
    assert (Path(updated.run_root) / "elements" / "E001" / "asset_package.json").exists()
    assert not (Path(updated.run_root) / "svg_to_ppt" / "assets" / "asset_manifest.json").exists()
    assert store.get_batch(batch.batch_id).status == "waiting_review"


def test_runner_auto_run_approves_and_reconstructs(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=True,
        config_path=base_config,
    )
    case_root = store.runs_root / batch.batch_id / "case_seed"
    config_path = create_case_config(
        base_config_path=base_config,
        source_image=source,
        output_dir=case_root,
        target_path=case_root / "drawai.config.yaml",
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=config_path,
    )
    runner = WorkbenchRunner(store, _settings(tmp_path, base_config), stage_executor=_deterministic_stage_executor)

    runner.submit_batch(batch.batch_id)
    runner.wait_for_idle(timeout=5)

    updated = store.get_case(case.case_id)
    assert updated.status == "completed"
    assert (Path(updated.run_root) / "reports" / "workbench" / "approved_asset_plan.json").exists()
    assert (Path(updated.run_root) / "svg" / "semantic.svg").exists()
    assert store.get_batch(batch.batch_id).status == "completed"


def test_runner_svg_rerun_executes_export_by_default(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=create_case_config(
            base_config_path=base_config,
            source_image=source,
            output_dir=store.runs_root / batch.batch_id / "case_seed",
            target_path=store.runs_root / batch.batch_id / "case_seed" / "drawai.config.yaml",
        ),
    )
    runner = WorkbenchRunner(store, _settings(tmp_path, base_config), stage_executor=_export_failing_stage_executor)

    runner.submit_rerun(case.case_id, "svg")
    runner.wait_for_idle(timeout=5)

    updated = store.get_case(case.case_id)
    labels = {artifact.label for artifact in store.list_artifacts(case.case_id)}
    stages = [stage.stage_name for stage in store.list_stage_runs(case.case_id)]
    assert updated.status == "failed"
    assert "RuntimeError: export unavailable" in (updated.error_message or "")
    assert stages == ["process_assets", "compose_svg", "export"]
    assert "semantic_svg" in labels
    assert "rendered_png" in labels
    assert "pptx" not in labels


def test_runner_archives_existing_svg_outputs_before_rerun(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=create_case_config(
            base_config_path=base_config,
            source_image=source,
            output_dir=store.runs_root / batch.batch_id / "case_seed",
            target_path=store.runs_root / batch.batch_id / "case_seed" / "drawai.config.yaml",
        ),
    )
    root = Path(case.run_root)
    (root / "svg").mkdir(parents=True, exist_ok=True)
    (root / "svg" / "semantic.svg").write_text("<svg>old</svg>\n", encoding="utf-8")
    (root / "svg" / "attempts" / "codex_merged" / "001").mkdir(parents=True, exist_ok=True)
    (root / "svg" / "attempts" / "codex_merged" / "001" / "model_response.txt").write_text(
        "old response\n",
        encoding="utf-8",
    )
    transient_browser_dir = root / "svg" / "attempts" / "codex_merged" / "001" / "chrome-profile-test"
    transient_browser_dir.mkdir(parents=True, exist_ok=True)
    (transient_browser_dir / "lock").write_text("volatile browser cache\n", encoding="utf-8")
    Image.new("RGB", (24, 24), "white").save(root / "svg" / "rendered.png")
    _write_json(root / "reports" / "svg_validation_report.json", {"status": "old"})

    def new_svg_executor(case_record, stage: str) -> None:
        if stage != "compose_svg":
            _deterministic_stage_executor(case_record, stage)
            return
        case_root = Path(case_record.run_root)
        (case_root / "svg").mkdir(parents=True, exist_ok=True)
        (case_root / "svg" / "semantic.svg").write_text("<svg>new</svg>\n", encoding="utf-8")
        Image.new("RGB", (24, 24), "black").save(case_root / "svg" / "rendered.png")
        _write_json(case_root / "reports" / "svg_validation_report.json", {"status": "new"})

    runner = WorkbenchRunner(store, _settings(tmp_path, base_config), stage_executor=new_svg_executor)

    runner.submit_rerun(case.case_id, "svg")
    runner.wait_for_idle(timeout=5)

    archives = sorted((root / "archives" / "svg_runs").glob("*_before_svg_rerun*"))
    assert len(archives) == 1
    assert (archives[0] / "svg" / "semantic.svg").read_text(encoding="utf-8") == "<svg>old</svg>\n"
    archived_response = archives[0] / "svg" / "attempts" / "codex_merged" / "001" / "model_response.txt"
    archived_transient = archives[0] / "svg" / "attempts" / "codex_merged" / "001" / "chrome-profile-test"
    assert Path(_archive_fs_path(archived_response)).exists()
    assert not Path(_archive_fs_path(archived_transient)).exists()
    assert (root / "svg" / "semantic.svg").read_text(encoding="utf-8") == "<svg>new</svg>\n"
    archive_manifest = json.loads((archives[0] / "archive_manifest.json").read_text(encoding="utf-8"))
    assert all("chrome-profile" not in path for path in archive_manifest["files"])
    labels = {artifact.label for artifact in store.list_artifacts(case.case_id)}
    assert "svg_rerun_archive" in labels


def test_runner_rejects_duplicate_case_jobs(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=create_case_config(
            base_config_path=base_config,
            source_image=source,
            output_dir=store.runs_root / batch.batch_id / "case_seed",
            target_path=store.runs_root / batch.batch_id / "case_seed" / "drawai.config.yaml",
        ),
    )
    started = threading.Event()
    release = threading.Event()

    def blocking_stage_executor(case_record, stage: str) -> None:
        if stage == "compose_svg":
            started.set()
            release.wait(timeout=5)
        _deterministic_stage_executor(case_record, stage)

    runner = WorkbenchRunner(store, _settings(tmp_path, base_config), stage_executor=blocking_stage_executor)
    runner.submit_rerun(case.case_id, "svg")
    assert started.wait(timeout=5)

    with pytest.raises(RuntimeError, match="active background job"):
        runner.submit_rerun(case.case_id, "svg")

    release.set()
    runner.wait_for_idle(timeout=5)


def test_runner_reports_resource_queue_activity(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    settings = WorkbenchSettings(
        workspace=tmp_path / "workspace",
        default_config=base_config,
        max_concurrent_cases=2,
        codex_concurrency=1,
    )
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="local_dir",
        max_concurrent_cases=2,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    for index in range(2):
        case_root = store.runs_root / batch.batch_id / f"case_{index}"
        store.create_case(
            batch_id=batch.batch_id,
            name=f"source_{index}.png",
            source_image_path=source,
            config_path=create_case_config(
                base_config_path=base_config,
                source_image=source,
                output_dir=case_root,
                target_path=case_root / "drawai.config.yaml",
            ),
        )
    first_codex_started = threading.Event()
    release = threading.Event()

    def blocking_codex_executor(case_record, stage: str) -> None:
        if stage == "refine_elements" and not first_codex_started.is_set():
            first_codex_started.set()
            release.wait(timeout=5)
        _deterministic_stage_executor(case_record, stage)

    runner = WorkbenchRunner(store, settings, stage_executor=blocking_codex_executor)

    runner.submit_batch(batch.batch_id)
    assert first_codex_started.wait(timeout=5)
    assert _wait_for_resource_activity(runner, "codex", queued=1, running=1, timeout=5)

    release.set()
    runner.wait_for_idle(timeout=5)
    codex_activity = runner.resource_activity()["codex"]
    assert codex_activity["queued"] == 0
    assert codex_activity["running"] == 0


def test_runner_parse_elements_uses_ocr_resource_lane(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    settings = WorkbenchSettings(
        workspace=tmp_path / "workspace",
        default_config=base_config,
        max_concurrent_cases=2,
        sam_concurrency=2,
        ocr_concurrency=1,
    )
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="local_dir",
        max_concurrent_cases=2,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    for index in range(2):
        case_root = store.runs_root / batch.batch_id / f"case_{index}"
        store.create_case(
            batch_id=batch.batch_id,
            name=f"source_{index}.png",
            source_image_path=source,
            config_path=create_case_config(
                base_config_path=base_config,
                source_image=source,
                output_dir=case_root,
                target_path=case_root / "drawai.config.yaml",
            ),
        )
    first_parse_started = threading.Event()
    release = threading.Event()

    def blocking_parse_executor(case_record, stage: str) -> None:
        if stage == "parse_elements" and not first_parse_started.is_set():
            first_parse_started.set()
            release.wait(timeout=5)
        _deterministic_stage_executor(case_record, stage)

    runner = WorkbenchRunner(store, settings, stage_executor=blocking_parse_executor)

    runner.submit_batch(batch.batch_id)
    assert first_parse_started.wait(timeout=5)
    assert _wait_for_resource_activity(runner, "ocr", queued=1, running=1, timeout=5)

    release.set()
    runner.wait_for_idle(timeout=5)
    ocr_activity = runner.resource_activity()["ocr"]
    assert ocr_activity["queued"] == 0
    assert ocr_activity["running"] == 0


def test_runner_marks_interrupted_running_case_failed_on_startup(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )
    stage_run = store.start_stage_run(case.case_id, "svg")
    store.update_case_status(case.case_id, status="svg_running", phase="reconstruction", stage="svg")
    store.update_batch_status(batch.batch_id, "running")

    WorkbenchRunner(store, _settings(tmp_path, base_config))

    updated_case = store.get_case(case.case_id)
    updated_stage = store.list_stage_runs(case.case_id)[0]
    assert updated_case.status == "failed"
    assert "interrupted" in updated_case.error_message
    assert updated_stage.stage_run_id == stage_run.stage_run_id
    assert updated_stage.status == "failed"
    assert "interrupted" in updated_stage.error_message
    assert store.get_batch(batch.batch_id).status == "failed"


def test_runner_marks_interrupted_running_stage_run_failed_on_startup(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )
    stage_run = store.start_stage_run(case.case_id, "svg")
    store.update_case_status(
        case.case_id,
        status="failed",
        phase="reconstruction",
        stage="svg",
        error_message="RuntimeError: failed",
    )
    store.update_batch_status(batch.batch_id, "failed", error_message="RuntimeError: failed")

    WorkbenchRunner(store, _settings(tmp_path, base_config))

    updated_case = store.get_case(case.case_id)
    updated_stage = store.list_stage_runs(case.case_id)[0]
    assert updated_case.status == "failed"
    assert updated_case.error_message == "RuntimeError: failed"
    assert updated_stage.stage_run_id == stage_run.stage_run_id
    assert updated_stage.status == "failed"
    assert "interrupted" in updated_stage.error_message


def test_runner_wait_for_idle_refreshes_stale_running_batch(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="upload",
        max_concurrent_cases=2,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    for index in range(2):
        case_root = store.runs_root / batch.batch_id / f"case_{index}"
        case = store.create_case(
            batch_id=batch.batch_id,
            name=f"source_{index}.png",
            source_image_path=source,
            config_path=create_case_config(
                base_config_path=base_config,
                source_image=source,
                output_dir=case_root,
                target_path=case_root / "drawai.config.yaml",
            ),
        )
        store.update_case_status(
            case.case_id,
            status="assets_review",
            phase="analysis",
            stage="plan_assets",
        )
    store.update_batch_status(batch.batch_id, "running")
    runner = WorkbenchRunner(store, _settings(tmp_path, base_config), stage_executor=_deterministic_stage_executor)

    runner.wait_for_idle(timeout=1)

    assert store.get_batch(batch.batch_id).status == "waiting_review"


def test_runner_copies_case_failure_to_batch_error(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="local_dir",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case_root = store.runs_root / batch.batch_id / "case_seed"
    config_path = create_case_config(
        base_config_path=base_config,
        source_image=source,
        output_dir=case_root,
        target_path=case_root / "drawai.config.yaml",
    )
    store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=config_path,
    )
    runner = WorkbenchRunner(store, _settings(tmp_path, base_config), stage_executor=_failing_stage_executor)

    runner.submit_batch(batch.batch_id)
    runner.wait_for_idle(timeout=5)

    updated = store.get_batch(batch.batch_id)
    assert updated.status == "failed"
    assert "detector unavailable" in updated.error_message


def test_api_exposes_v2_package_and_asset_package(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    settings = _settings(tmp_path, base_config)
    runner = WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor)
    app = create_app(settings, store=store, runner=runner)
    client = TestClient(app)
    batch = store.create_batch(
        name="v2 batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )
    _write_minimal_v2_package(Path(case.run_root), case.case_id)

    package_response = client.get(f"/api/cases/{case.case_id}/package")
    elements_response = client.get(f"/api/cases/{case.case_id}/elements")
    asset_response = client.get(f"/api/cases/{case.case_id}/elements/E001/asset-package")
    case_response = client.get(f"/api/cases/{case.case_id}")

    assert package_response.status_code == 200
    assert package_response.json()["package"]["schema"] == "drawai.run_package.v1"
    assert package_response.json()["compatibility"]["mode"] == "v2"
    assert elements_response.status_code == 200
    assert elements_response.json()["elements"][0]["element_id"] == "E001"
    assert asset_response.status_code == 200
    assert asset_response.json()["asset_package"]["element_id"] == "E001"
    assert case_response.status_code == 200
    assert case_response.json()["case"]["compatibility_mode"] == "v2"
    assert case_response.json()["case"]["can_fork_from_source"] is True


def test_api_asset_process_marks_downstream_outputs_stale(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    settings = _settings(tmp_path, base_config)
    runner = WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor)
    app = create_app(settings, store=store, runner=runner)
    client = TestClient(app)
    batch = store.create_batch(
        name="v2 batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )
    _write_minimal_v2_package(Path(case.run_root), case.case_id)
    package_path = Path(case.run_root) / "drawai_package.json"
    package = json.loads(package_path.read_text(encoding="utf-8"))
    package["compose_outputs"] = {"semantic_svg": "svg/semantic.svg"}
    package["export_outputs"] = {"report": "reports/svg_to_ppt_export_report.json"}
    package_path.write_text(json.dumps(package, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    store.update_case_status(case.case_id, status="completed", phase="reconstruction", stage="completed")

    response = client.post(
        f"/api/cases/{case.case_id}/elements/E001/process",
        json={"processor": "crop"},
    )

    assert response.status_code == 200
    updated_package = json.loads(package_path.read_text(encoding="utf-8"))
    assert "compose_outputs" not in updated_package
    assert "export_outputs" not in updated_package
    updated_case = store.get_case(case.case_id)
    assert updated_case.status == "assets_review"
    assert updated_case.stale_from_stage == "compose_svg"


def test_api_failed_asset_process_marks_downstream_outputs_stale(tmp_path: Path) -> None:
    class FailingRmbgClient:
        def remove_background(self, *_args: object, **_kwargs: object) -> RmbgResult:
            raise RuntimeError("rmbg unavailable")

    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    settings = _settings(tmp_path, base_config)
    runner = WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor)
    app = create_app(settings, store=store, runner=runner, rmbg_client=FailingRmbgClient())
    client = TestClient(app)
    batch = store.create_batch(
        name="v2 batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )
    _write_minimal_v2_package(Path(case.run_root), case.case_id)
    package_path = Path(case.run_root) / "drawai_package.json"
    package = json.loads(package_path.read_text(encoding="utf-8"))
    package["compose_outputs"] = {"semantic_svg": "svg/semantic.svg"}
    package["export_outputs"] = {"report": "reports/svg_to_ppt_export_report.json"}
    package_path.write_text(json.dumps(package, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    store.update_case_status(case.case_id, status="completed", phase="reconstruction", stage="completed")

    response = client.post(
        f"/api/cases/{case.case_id}/elements/E001/process",
        json={"processor": "crop_nobg"},
    )

    assert response.status_code == 400
    updated_package = json.loads(package_path.read_text(encoding="utf-8"))
    assert updated_package["asset_packages"][0]["status"] == "failed"
    assert "compose_outputs" not in updated_package
    assert "export_outputs" not in updated_package
    updated_case = store.get_case(case.case_id)
    assert updated_case.status == "assets_review"
    assert updated_case.stale_from_stage == "compose_svg"


def test_api_asset_process_file_error_marks_downstream_outputs_stale(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    settings = _settings(tmp_path, base_config)
    runner = WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor)
    app = create_app(settings, store=store, runner=runner)
    client = TestClient(app)
    batch = store.create_batch(
        name="v2 batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )
    _write_minimal_v2_package(Path(case.run_root), case.case_id)
    package_path = Path(case.run_root) / "drawai_package.json"
    package = json.loads(package_path.read_text(encoding="utf-8"))
    package["source_image"] = str(Path(case.run_root) / "inputs" / "missing.png")
    package["compose_outputs"] = {"semantic_svg": "svg/semantic.svg"}
    package["export_outputs"] = {"report": "reports/svg_to_ppt_export_report.json"}
    package_path.write_text(json.dumps(package, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    store.update_case_status(case.case_id, status="completed", phase="reconstruction", stage="completed")

    response = client.post(
        f"/api/cases/{case.case_id}/elements/E001/process",
        json={"processor": "crop"},
    )

    assert response.status_code == 400
    updated_package = json.loads(package_path.read_text(encoding="utf-8"))
    assert updated_package["asset_packages"][0]["status"] == "failed"
    assert "compose_outputs" not in updated_package
    assert "export_outputs" not in updated_package
    updated_case = store.get_case(case.case_id)
    assert updated_case.status == "assets_review"
    assert updated_case.stale_from_stage == "compose_svg"


def test_legacy_case_mutation_is_rejected_but_can_fork_from_source(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    settings = _settings(tmp_path, base_config)
    runner = WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor)
    app = create_app(settings, store=store, runner=runner)
    client = TestClient(app)
    batch = store.create_batch(
        name="legacy batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )
    root = Path(case.run_root)
    (root / "inputs").mkdir(parents=True)
    shutil.copy2(source, root / "inputs" / "figure.png")
    (root / "svg").mkdir()
    (root / "svg" / "semantic.svg").write_text("<svg />\n", encoding="utf-8")

    process_response = client.post(
        f"/api/cases/{case.case_id}/elements/E001/process",
        json={"processor": "crop"},
    )
    retry_response = client.post(f"/api/cases/{case.case_id}/retry")
    fork_response = client.post(f"/api/cases/{case.case_id}/fork-v2-from-source")

    assert process_response.status_code == 409
    assert process_response.json()["detail"] == "legacy_readonly_case"
    assert retry_response.status_code == 409
    assert retry_response.json()["detail"] == "legacy_readonly_case"
    assert fork_response.status_code == 200
    forked_case = fork_response.json()["case"]
    assert forked_case["case_id"] != case.case_id
    runner.wait_for_idle(timeout=5)
    assert (Path(forked_case["run_root"]) / "drawai_package.json").exists()


def test_api_creates_batch_polls_assets_and_approves(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    image_dir = tmp_path / "images"
    image_dir.mkdir()
    Image.new("RGB", (24, 24), "white").save(image_dir / "source.png")
    settings = _settings(tmp_path, base_config)
    runner = WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor)
    app = create_app(settings, store=store, runner=runner)
    client = TestClient(app)

    response = client.post(
        "/api/batches",
        json={
            "name": "api batch",
            "input_mode": "local_dir",
            "local_dir": str(image_dir),
            "auto_run_svg_after_analysis": False,
            "max_concurrent_cases": 1,
            "base_config_path": str(base_config),
        },
    )
    assert response.status_code == 200
    batch_id = response.json()["batch"]["batch_id"]
    case_id = response.json()["cases"][0]["case_id"]

    runner.wait_for_idle(timeout=5)
    batch_response = client.get(f"/api/batches/{batch_id}")
    assert batch_response.json()["batch"]["status"] == "waiting_review"
    assert batch_response.json()["cases"][0]["preview_url"] == f"/api/cases/{case_id}/source-image"
    assert batch_response.json()["cases"][0]["editor_ready"] is True
    assets_response = client.get(f"/api/cases/{case_id}/assets")
    assert assets_response.status_code == 200
    plan = assets_response.json()["asset_plan"]
    plan["elements"][0]["source_strategy"] = "svg_self_draw"

    save_response = client.patch(f"/api/cases/{case_id}/asset-draft", json=plan)
    assert save_response.status_code == 200
    approve_response = client.post(f"/api/cases/{case_id}/approve-assets", json={"run_svg": False})
    assert approve_response.status_code == 200
    case_response = client.get(f"/api/cases/{case_id}")
    case_payload = case_response.json()["case"]
    labels = {artifact["label"] for artifact in case_response.json()["artifacts"]}
    assert case_payload["status"] == "assets_review"
    assert case_payload["stage"] == "approved_asset_plan"
    assert "approved_asset_plan" in labels
    assert "asset_manifest" not in labels
    assert not (Path(case_payload["run_root"]) / "svg_to_ppt" / "assets" / "asset_manifest.json").exists()


def test_api_approve_with_run_svg_enters_running_state_immediately(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    image_dir = tmp_path / "images"
    image_dir.mkdir()
    Image.new("RGB", (24, 24), "white").save(image_dir / "source.png")
    settings = _settings(tmp_path, base_config)
    started = threading.Event()
    release = threading.Event()

    def blocking_process_assets(case_record, stage: str) -> None:
        if stage == "process_assets":
            started.set()
            release.wait(timeout=5)
        _deterministic_stage_executor(case_record, stage)

    runner = WorkbenchRunner(store, settings, stage_executor=blocking_process_assets)
    app = create_app(settings, store=store, runner=runner)
    client = TestClient(app)

    response = client.post(
        "/api/batches",
        json={
            "name": "api batch",
            "input_mode": "local_dir",
            "local_dir": str(image_dir),
            "auto_run_svg_after_analysis": False,
            "max_concurrent_cases": 1,
            "base_config_path": str(base_config),
        },
    )
    assert response.status_code == 200
    case_id = response.json()["cases"][0]["case_id"]
    runner.wait_for_idle(timeout=5)

    approve_response = client.post(f"/api/cases/{case_id}/approve-assets", json={"run_svg": True})

    assert approve_response.status_code == 200
    case_payload = approve_response.json()["case"]
    assert case_payload["status"] == "svg_running"
    assert case_payload["stage"] == "process_assets"
    assert started.wait(timeout=5)
    assert store.get_batch(case_payload["batch_id"]).status == "running"

    release.set()
    runner.wait_for_idle(timeout=5)
    updated = store.get_case(case_id)
    assert updated.status == "completed"
    assert (Path(updated.run_root) / "svg" / "semantic.svg").exists()


def test_api_renames_and_deletes_batch(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    settings = _settings(tmp_path, base_config)
    app = create_app(settings, store=store, runner=WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor))
    client = TestClient(app)
    batch = store.create_batch(
        name="old task",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )
    upload_root = store.uploads_root / batch.batch_id
    upload_root.mkdir(parents=True)
    (upload_root / "source.png").write_bytes(b"upload")
    (Path(case.run_root) / "scratch.txt").write_text("run file", encoding="utf-8")

    rename_response = client.patch(f"/api/batches/{batch.batch_id}", json={"name": "renamed task"})

    assert rename_response.status_code == 200
    assert rename_response.json()["batch"]["name"] == "renamed task"
    assert store.get_batch(batch.batch_id).name == "renamed task"

    delete_response = client.delete(f"/api/batches/{batch.batch_id}")

    assert delete_response.status_code == 200
    assert client.get(f"/api/batches/{batch.batch_id}").status_code == 404
    assert store.list_cases(batch.batch_id) == []
    assert not (store.runs_root / batch.batch_id).exists()
    assert not upload_root.exists()


def test_api_run_batch_approves_asset_drafts_and_generates_svg(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="task",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case_root = store.runs_root / batch.batch_id / "case_seed"
    config_path = create_case_config(
        base_config_path=base_config,
        source_image=source,
        output_dir=case_root,
        target_path=case_root / "drawai.config.yaml",
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=config_path,
    )
    settings = _settings(tmp_path, base_config)
    runner = WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor)
    app = create_app(settings, store=store, runner=runner)
    client = TestClient(app)
    runner.submit_batch(batch.batch_id)
    runner.wait_for_idle(timeout=5)
    assert store.get_case(case.case_id).status == "assets_review"

    response = client.post(f"/api/batches/{batch.batch_id}/run")

    assert response.status_code == 200
    runner.wait_for_idle(timeout=5)
    updated = store.get_case(case.case_id)
    assert updated.status == "completed"
    assert store.get_batch(batch.batch_id).status == "completed"
    assert (Path(updated.run_root) / "svg" / "semantic.svg").exists()


def test_api_downloads_completed_batch_as_merged_pptx(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    settings = _settings(tmp_path, base_config)
    app = create_app(settings, store=store, runner=WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor))
    client = TestClient(app)
    batch = store.create_batch(
        name="paper figures",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    first = store.create_case(batch_id=batch.batch_id, name="first.png", source_image_path=source, config_path=base_config)
    second = store.create_case(batch_id=batch.batch_id, name="second.png", source_image_path=source, config_path=base_config)
    for case, text in [(first, "First slide"), (second, "Second slide")]:
        pptx_path = Path(case.run_root) / "svg_to_ppt" / "semantic.svg_to_ppt.pptx"
        _write_single_text_slide_pptx(pptx_path, text)
        store.register_artifact(case.case_id, label="pptx", path=pptx_path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        store.update_case_status(case.case_id, status="completed", phase="reconstruction", stage="completed")
    store.update_batch_status(batch.batch_id, "completed")

    response = client.get(f"/api/batches/{batch.batch_id}/pptx")

    assert response.status_code == 200
    assert response.headers["content-type"] == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    merged = tmp_path / "merged.pptx"
    merged.write_bytes(response.content)
    presentation = Presentation(str(merged))
    assert len(presentation.slides) == 2
    slide_texts = [" ".join(shape.text for shape in slide.shapes if hasattr(shape, "text")) for slide in presentation.slides]
    assert "First slide" in slide_texts[0]
    assert "Second slide" in slide_texts[1]


def test_api_creates_batch_from_uploaded_zip(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    settings = WorkbenchSettings(workspace=tmp_path / "workspace", default_config=base_config)
    runner = WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor)
    app = create_app(settings, store=store, runner=runner)
    client = TestClient(app)
    image_bytes = io.BytesIO()
    Image.new("RGB", (24, 24), "white").save(image_bytes, format="PNG")
    archive_bytes = io.BytesIO()
    with zipfile.ZipFile(archive_bytes, "w") as archive:
        archive.writestr("figures/source.png", image_bytes.getvalue())
        archive.writestr("notes/readme.txt", "ignored")
    archive_bytes.seek(0)

    response = client.post(
        "/api/batches",
        data={
            "name": "zip upload batch",
            "input_mode": "upload",
        },
        files={"files": ("figures.zip", archive_bytes.getvalue(), "application/zip")},
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["batch"]["input_mode"] == "upload"
    assert payload["batch"]["max_concurrent_cases"] == 10
    assert payload["batch"]["config_path"] == str(base_config)
    assert len(payload["cases"]) == 1
    assert payload["cases"][0]["name"] == "source.png"
    runner.wait_for_idle(timeout=5)
    assert client.get(f"/api/batches/{payload['batch']['batch_id']}").json()["batch"]["status"] == "waiting_review"


def test_api_creates_batch_from_multiple_uploaded_images(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    settings = WorkbenchSettings(workspace=tmp_path / "workspace", default_config=base_config)
    runner = WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor)
    app = create_app(settings, store=store, runner=runner)
    client = TestClient(app)
    first = io.BytesIO()
    second = io.BytesIO()
    Image.new("RGB", (24, 24), "white").save(first, format="PNG")
    Image.new("RGB", (32, 20), "black").save(second, format="PNG")

    response = client.post(
        "/api/batches",
        data={
            "name": "multi image upload",
            "input_mode": "upload",
        },
        files=[
            ("files", ("alpha.png", first.getvalue(), "image/png")),
            ("files", ("nested/beta.png", second.getvalue(), "image/png")),
        ],
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["batch"]["input_mode"] == "upload"
    assert payload["batch"]["case_counts"] == {"queued": 2}
    assert len(payload["cases"]) == 2
    assert [case["name"] for case in payload["cases"]] == ["alpha.png", "beta.png"]
    runner.wait_for_idle(timeout=5)
    batch_payload = client.get(f"/api/batches/{payload['batch']['batch_id']}").json()
    assert batch_payload["batch"]["status"] == "waiting_review"
    assert batch_payload["batch"]["case_counts"]["assets_review"] == 2


def test_api_creates_batch_from_generated_image_data_url(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    settings = WorkbenchSettings(workspace=tmp_path / "workspace", default_config=base_config)
    runner = WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor)
    app = create_app(settings, store=store, runner=runner)
    client = TestClient(app)
    image_bytes = io.BytesIO()
    Image.new("RGB", (28, 18), "white").save(image_bytes, format="PNG")
    data_url = f"data:image/png;base64,{base64.b64encode(image_bytes.getvalue()).decode('ascii')}"

    response = client.post(
        "/api/batches",
        data={
            "name": "generated upload",
            "input_mode": "upload",
            "generated_image_urls": data_url,
        },
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["batch"]["input_mode"] == "upload"
    assert len(payload["cases"]) == 1
    assert payload["cases"][0]["name"] == "generated-001.png"
    runner.wait_for_idle(timeout=5)
    batch_payload = client.get(f"/api/batches/{payload['batch']['batch_id']}").json()
    assert batch_payload["batch"]["status"] == "waiting_review"
    assert batch_payload["batch"]["case_counts"]["assets_review"] == 1


def test_api_processes_asset_and_returns_model_output_url(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (16, 16), "white").save(source)
    settings = _settings(tmp_path, base_config)
    runner = WorkbenchRunner(store, settings)
    app = create_app(settings, store=store, runner=runner, rmbg_client=FakeApiRmbgClient())
    client = TestClient(app)
    batch = store.create_batch(
        name="processing batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )
    figure = Image.new("RGBA", (16, 16), (255, 255, 255, 255))
    for x in range(4, 12):
        for y in range(4, 12):
            figure.putpixel((x, y), (10, 160, 90, 255))
    figure_path = Path(case.run_root) / "inputs" / "figure.png"
    figure_path.parent.mkdir(parents=True, exist_ok=True)
    figure.save(figure_path)
    plan = {
        "schema": "drawai.workbench_asset_plan.v1",
        "case_id": case.case_id,
        "elements": [
            {
                "box_id": "A001",
                "bbox": [2, 2, 14, 14],
                "source_strategy": "crop_nobg",
                "type": "image",
            }
        ],
    }

    response = client.post(
        f"/api/cases/{case.case_id}/asset-processing",
        json={"asset_plan": plan, "asset_ids": ["A001"]},
    )

    assert response.status_code == 200
    payload = response.json()
    element = payload["asset_plan"]["elements"][0]
    assert element["processing_status"] == "processed"
    assert element["processed_asset_source_strategy"] == "crop_nobg"
    assert element["processed_asset_relative_path"].endswith("A001_nobg.png")
    image_response = client.get(payload["processed_assets"][0]["url"])
    assert image_response.status_code == 200
    with Image.open(io.BytesIO(image_response.content)) as image:
        assert image.mode == "RGBA"
        assert image.getpixel((0, 0))[3] == 0


def test_api_writes_model_runtime_urls_into_case_config(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    image_dir = tmp_path / "images"
    image_dir.mkdir()
    Image.new("RGB", (24, 24), "white").save(image_dir / "source.png")
    settings = WorkbenchSettings(
        workspace=tmp_path / "workspace",
        default_config=base_config,
        sam3_base_url="http://model-a:18080",
        ocr_base_url="http://model-a:18080",
        rmbg_base_url="http://model-a:18080",
        ocr_timeout_seconds=600,
    )
    runner = WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor)
    app = create_app(settings, store=store, runner=runner)
    client = TestClient(app)

    response = client.post(
        "/api/batches",
        json={
            "name": "remote model batch",
            "input_mode": "local_dir",
            "local_dir": str(image_dir),
            "auto_run_svg_after_analysis": False,
            "max_concurrent_cases": 1,
            "base_config_path": str(base_config),
        },
    )

    assert response.status_code == 200
    case_id = response.json()["cases"][0]["case_id"]
    case = store.get_case(case_id)
    payload = yaml.safe_load(Path(case.config_path).read_text(encoding="utf-8"))
    assert payload["sam3"]["base_url"] == "http://model-a:18080"
    assert payload["ocr"]["remote_paddleocr"]["base_url"] == "http://model-a:18080"
    assert payload["ocr"]["remote_paddleocr"]["timeout_seconds"] == 600
    assert payload["asset_materialization"]["rmbg"]["base_url"] == "http://model-a:18080"


def test_settings_from_env_defaults_codex_concurrency_to_five(monkeypatch, tmp_path: Path) -> None:
    monkeypatch.setenv("DRAWAI_WORKBENCH_WORKSPACE", str(tmp_path / "workspace"))
    monkeypatch.setenv("DRAWAI_WORKBENCH_DEFAULT_CONFIG", str(tmp_path / "config.yaml"))
    monkeypatch.delenv("DRAWAI_WORKBENCH_CODEX_CONCURRENCY", raising=False)

    settings = workbench_api.settings_from_env()

    assert settings.codex_concurrency == 5

    monkeypatch.setenv("DRAWAI_WORKBENCH_CODEX_CONCURRENCY", "2")

    settings = workbench_api.settings_from_env()

    assert settings.codex_concurrency == 2


def test_rerun_refreshes_case_runtime_config(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case_root = store.runs_root / batch.batch_id / "case_seed"
    stale_config = create_case_config(
        base_config_path=base_config,
        source_image=source,
        output_dir=case_root,
        target_path=case_root / "drawai.config.yaml",
        ocr_base_url="http://old-model:18080",
        ocr_timeout_seconds=60,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=stale_config,
    )
    settings = WorkbenchSettings(
        workspace=tmp_path / "workspace",
        default_config=base_config,
        ocr_base_url="http://model-a:18080",
        ocr_timeout_seconds=600,
    )
    runner = WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor)

    runner.submit_rerun(case.case_id, "analysis")
    runner.wait_for_idle(timeout=5)

    payload = yaml.safe_load(Path(store.get_case(case.case_id).config_path).read_text(encoding="utf-8"))
    assert payload["ocr"]["remote_paddleocr"]["base_url"] == "http://model-a:18080"
    assert payload["ocr"]["remote_paddleocr"]["timeout_seconds"] == 600


def test_api_run_stage_accepts_v2_boundary_stage_names(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    settings = _settings(tmp_path, base_config)
    runner = WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor)
    app = create_app(settings, store=store, runner=runner)
    client = TestClient(app)
    batch = store.create_batch(
        name="v2 retry names",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )

    prepare_response = client.post(f"/api/cases/{case.case_id}/run-stage", json={"stage": "prepare"})
    runner.wait_for_idle(timeout=5)
    package_response = client.post(f"/api/cases/{case.case_id}/run-stage", json={"stage": "package_run"})
    runner.wait_for_idle(timeout=5)

    assert prepare_response.status_code == 200
    assert package_response.status_code == 200
    assert store.list_stage_runs(case.case_id)[0].stage_name == "prepare"
    assert any(stage.stage_name == "export" for stage in store.list_stage_runs(case.case_id))


def test_api_case_list_uses_source_image_preview_before_artifacts(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    settings = _settings(tmp_path, base_config)
    runner = WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor)
    app = create_app(settings, store=store, runner=runner)
    client = TestClient(app)
    batch = store.create_batch(
        name="pending batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )
    _deterministic_stage_executor(case, "prepare")
    _write_minimal_v2_package(Path(case.run_root), case.case_id)
    _deterministic_stage_executor(case, "compose_svg")
    store.register_artifact(case.case_id, label="figure", path=Path(case.run_root) / "inputs" / "figure.png", media_type="image/png")
    store.register_artifact(case.case_id, label="rendered_png", path=Path(case.run_root) / "svg" / "rendered.png", media_type="image/png")

    batch_response = client.get(f"/api/batches/{batch.batch_id}")

    assert batch_response.status_code == 200
    preview_url = batch_response.json()["cases"][0]["preview_url"]
    assert preview_url == f"/api/cases/{case.case_id}/source-image"
    image_response = client.get(preview_url)
    assert image_response.status_code == 200
    assert image_response.headers["content-type"] == "image/png"


def test_api_health_reports_runtime_services(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    settings = _settings(tmp_path, base_config)
    app = create_app(
        settings,
        store=store,
        runner=WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor),
        runtime_probe=lambda name, base_url: {
            "name": name,
            "base_url": base_url,
            "health_url": f"{base_url}/health",
            "status": "online",
        },
    )
    client = TestClient(app)

    response = client.get("/api/health")

    assert response.status_code == 200
    payload = response.json()
    assert payload["status"] == "ok"
    runtime_services = payload["runtime_services"]
    assert runtime_services["sam3"]["status"] == "online"
    assert runtime_services["ocr"]["base_url"] == "http://127.0.0.1:18080"
    assert runtime_services["rmbg"]["base_url"] == "http://127.0.0.1:18080"
    assert payload["runtime_activity"]["sam3"] == {"limit": 1, "queued": 0, "running": 0}
    assert payload["runtime_activity"]["codex"] == {"limit": 5, "queued": 0, "running": 0}


def test_api_health_is_degraded_when_any_runtime_service_is_offline(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    settings = _settings(tmp_path, base_config)

    def probe(name: str, base_url: str) -> dict[str, str]:
        status = "offline" if name == "ocr" else "online"
        payload = {
            "name": name,
            "base_url": base_url,
            "health_url": f"{base_url}/health",
            "status": status,
        }
        if status == "offline":
            payload["error"] = "connection refused"
        return payload

    app = create_app(
        settings,
        store=store,
        runner=WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor),
        runtime_probe=probe,
    )
    client = TestClient(app)

    response = client.get("/api/health")

    assert response.status_code == 200
    payload = response.json()
    assert payload["status"] == "degraded"
    assert payload["runtime_services"]["ocr"]["status"] == "offline"


def test_api_case_progress_exposes_case_files(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )
    _write_minimal_v2_package(Path(case.run_root), case.case_id)
    _deterministic_stage_executor(case, "compose_svg")
    app = create_app(_settings(tmp_path, base_config), store=store, runner=WorkbenchRunner(store, _settings(tmp_path, base_config)))
    client = TestClient(app)

    progress = client.get(f"/api/cases/{case.case_id}/progress")

    assert progress.status_code == 200
    semantic = next(item for item in progress.json()["files"] if item["label"] == "semantic_svg")
    assert semantic["exists"] is True
    assert semantic["url"].endswith("/svg/semantic.svg")
    file_response = client.get(semantic["url"])
    assert file_response.status_code == 200
    assert file_response.text.startswith("<svg")


def test_api_case_progress_exposes_pptx_export_mode(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )
    _deterministic_stage_executor(case, "export")
    app = create_app(_settings(tmp_path, base_config), store=store, runner=WorkbenchRunner(store, _settings(tmp_path, base_config)))
    client = TestClient(app)

    progress = client.get(f"/api/cases/{case.case_id}/progress")

    assert progress.status_code == 200
    pptx_export = progress.json()["pptx_export"]
    assert pptx_export["status"] == "ok"
    assert pptx_export["export_backend"] == "drawai_native_shapes"
    assert pptx_export["effective_export_mode"] == "native_shapes"
    assert pptx_export["editable_surface"] == "native_shapes"
    assert pptx_export["report_url"].endswith("/reports/svg_to_ppt_export_report.json")


def test_api_reads_and_updates_svg_source(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )
    _write_minimal_v2_package(Path(case.run_root), case.case_id)
    _deterministic_stage_executor(case, "compose_svg")
    package_path = Path(case.run_root) / "drawai_package.json"
    package = json.loads(package_path.read_text(encoding="utf-8"))
    package["compose_outputs"] = {"semantic_svg": "svg/semantic.svg"}
    package["export_outputs"] = {"report": "reports/svg_to_ppt_export_report.json"}
    package_path.write_text(json.dumps(package, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    store.update_case_status(case.case_id, status="completed", phase="reconstruction", stage="completed")
    store.update_batch_status(batch.batch_id, "completed")
    app = create_app(_settings(tmp_path, base_config), store=store, runner=WorkbenchRunner(store, _settings(tmp_path, base_config)))
    client = TestClient(app)

    read_response = client.get(f"/api/cases/{case.case_id}/svg-source")
    assert read_response.status_code == 200
    assert read_response.json()["svg"].startswith("<svg")

    edited_svg = (
        '<svg xmlns="http://www.w3.org/2000/svg" width="24" height="24" viewBox="0 0 24 24">'
        '<rect x="2" y="3" width="4" height="5"/></svg>'
    )
    update_response = client.patch(f"/api/cases/{case.case_id}/svg-source", json={"svg": edited_svg})

    assert update_response.status_code == 200
    assert (Path(case.run_root) / "svg" / "semantic.svg").read_text(encoding="utf-8").strip() == edited_svg
    updated_case = client.get(f"/api/cases/{case.case_id}").json()["case"]
    assert updated_case["stage"] == "svg_edit"
    assert updated_case["stale_from_stage"] == "export"
    updated_package = json.loads(package_path.read_text(encoding="utf-8"))
    assert updated_package["compose_outputs"] == {"semantic_svg": "svg/semantic.svg"}
    assert "export_outputs" not in updated_package
    labels = {artifact["label"] for artifact in client.get(f"/api/cases/{case.case_id}").json()["artifacts"]}
    assert "semantic_svg" in labels


def test_api_rejects_unsafe_svg_source(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )
    _write_minimal_v2_package(Path(case.run_root), case.case_id)
    _deterministic_stage_executor(case, "compose_svg")
    app = create_app(_settings(tmp_path, base_config), store=store, runner=WorkbenchRunner(store, _settings(tmp_path, base_config)))
    client = TestClient(app)

    response = client.patch(
        f"/api/cases/{case.case_id}/svg-source",
        json={"svg": '<svg xmlns="http://www.w3.org/2000/svg"><script>alert(1)</script></svg>'},
    )

    assert response.status_code == 400
    assert "active content" in response.json()["detail"]


def test_api_case_progress_summarizes_svg_attempt_failure(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "source.png"
    Image.new("RGB", (24, 24), "white").save(source)
    batch = store.create_batch(
        name="batch",
        input_mode="upload",
        max_concurrent_cases=1,
        auto_run_svg_after_analysis=False,
        config_path=base_config,
    )
    case = store.create_case(
        batch_id=batch.batch_id,
        name="source.png",
        source_image_path=source,
        config_path=base_config,
    )
    attempt_dir = Path(case.run_root) / "svg" / "attempts" / "codex_merged" / "001"
    (attempt_dir / "prompt.txt").parent.mkdir(parents=True, exist_ok=True)
    (attempt_dir / "prompt.txt").write_text("draw the figure", encoding="utf-8")
    _write_json(
        attempt_dir / "validation_report.json",
        {
            "status": "failed",
            "issues": [
                {"code": "model_text_missing_role", "message": "Text element is missing data-pb-role.", "detail": "✓"},
                {"code": "invoker_error", "message": "SVG generation invoker raised an exception.", "detail": "refresh token revoked"},
            ],
        },
    )
    app = create_app(_settings(tmp_path, base_config), store=store, runner=WorkbenchRunner(store, _settings(tmp_path, base_config)))
    client = TestClient(app)

    progress = client.get(f"/api/cases/{case.case_id}/progress")

    assert progress.status_code == 200
    attempt = progress.json()["svg_attempts"][0]
    assert attempt["phase"] == "codex_merged"
    assert attempt["attempt"] == "001"
    assert attempt["status"] == "failed"
    assert attempt["issue_count"] == 2
    assert "model_text_missing_role" in attempt["error_message"]
    assert "refresh token revoked" in attempt["error_message"]


def test_runtime_health_probe_recognizes_loopback_urls() -> None:
    assert is_loopback_url("http://127.0.0.1:18080/health")
    assert is_loopback_url("http://localhost:18080/health")
    assert is_loopback_url("http://[::1]:18080/health")
    assert not is_loopback_url("https://runtime.example.com/health")


def test_runtime_health_probe_marks_nested_service_failure_offline(monkeypatch: pytest.MonkeyPatch) -> None:
    payload = {
        "status": "ok",
        "services": {
            "sam3": {"status": "ok"},
            "ocr": {"status": "failed", "error": "model not loaded"},
            "rmbg": {"status": "ok"},
        },
    }

    class FakeResponse:
        def __enter__(self) -> "FakeResponse":
            return self

        def __exit__(self, exc_type: object, exc: object, tb: object) -> None:
            return None

        def read(self, _: int = -1) -> bytes:
            return json.dumps(payload).encode("utf-8")

    monkeypatch.setattr(workbench_api, "urlopen_direct_for_loopback", lambda *_args, **_kwargs: FakeResponse())

    status = workbench_api._probe_runtime_service("ocr", "http://127.0.0.1:18080")

    assert status["status"] == "offline"
    assert status["error"] == "model not loaded"


def test_image_generation_endpoint_requires_server_side_api_key(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.delenv("DRAWAI_IMAGEGEN_API_KEY", raising=False)
    monkeypatch.delenv("OPENAI_API_KEY", raising=False)
    base_config = _base_config(tmp_path)
    store = WorkbenchStore(tmp_path / "workspace")
    settings = _settings(tmp_path, base_config)
    app = create_app(settings, store=store, runner=WorkbenchRunner(store, settings))
    client = TestClient(app)

    response = client.post(
        "/api/imagegen/generations",
        json={"model": "gpt-image-2", "prompt": "draw a clean green logo", "n": 1},
    )

    assert response.status_code == 503
    assert "DRAWAI_IMAGEGEN_API_KEY" in response.json()["detail"]


def test_image_generation_endpoint_proxies_upstream_request(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setenv("DRAWAI_IMAGEGEN_API_KEY", "secret-test-key")
    monkeypatch.setenv("DRAWAI_IMAGEGEN_API_URL", "https://image.example.test/v1/images/generations")
    base_config = _base_config(tmp_path)
    store = WorkbenchStore(tmp_path / "workspace")
    settings = _settings(tmp_path, base_config)
    captured: dict[str, object] = {}

    class FakeResponse:
        def __enter__(self) -> "FakeResponse":
            return self

        def __exit__(self, exc_type: object, exc: object, tb: object) -> None:
            return None

        def read(self, _: int = -1) -> bytes:
            return json.dumps({"data": [{"url": "https://cdn.example.test/image.png"}]}).encode("utf-8")

    def fake_urlopen(request, *, timeout: float):
        captured["url"] = request.full_url
        captured["authorization"] = request.get_header("Authorization")
        captured["timeout"] = timeout
        captured["payload"] = json.loads(request.data.decode("utf-8"))
        return FakeResponse()

    monkeypatch.setattr(workbench_api, "urlopen_external", fake_urlopen)
    app = create_app(settings, store=store, runner=WorkbenchRunner(store, settings))
    client = TestClient(app)

    response = client.post(
        "/api/imagegen/generations",
        json={
            "model": "gpt-image-2",
            "prompt": "draw a clean green logo",
            "size": "1024x1024",
            "quality": "high",
            "background": "transparent",
            "moderation": "auto",
            "output_format": "png",
            "n": 1,
            "ignored": "not forwarded",
        },
    )

    assert response.status_code == 200
    assert response.json()["data"][0]["url"] == "https://cdn.example.test/image.png"
    assert captured["url"] == "https://image.example.test/v1/images/generations"
    assert captured["authorization"] == "Bearer secret-test-key"
    assert captured["timeout"] == 600.0
    payload = captured["payload"]
    assert isinstance(payload, dict)
    assert payload["model"] == "gpt-image-2"
    assert payload["prompt"] == "draw a clean green logo"
    assert payload["size"] == "1024x1024"
    assert payload["n"] == 1
    assert "resolution" not in payload
    assert "ignored" not in payload


def test_image_generation_endpoint_accepts_request_connection_settings(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.delenv("DRAWAI_IMAGEGEN_API_KEY", raising=False)
    monkeypatch.delenv("OPENAI_API_KEY", raising=False)
    monkeypatch.delenv("DRAWAI_IMAGEGEN_API_URL", raising=False)
    base_config = _base_config(tmp_path)
    store = WorkbenchStore(tmp_path / "workspace")
    settings = _settings(tmp_path, base_config)
    captured: dict[str, object] = {}

    class FakeResponse:
        def __enter__(self) -> "FakeResponse":
            return self

        def __exit__(self, exc_type: object, exc: object, tb: object) -> None:
            return None

        def read(self, _: int = -1) -> bytes:
            return json.dumps({"data": [{"url": "https://cdn.example.test/generated.png"}]}).encode("utf-8")

    def fake_urlopen(request, *, timeout: float):
        captured["url"] = request.full_url
        captured["authorization"] = request.get_header("Authorization")
        captured["payload"] = json.loads(request.data.decode("utf-8"))
        return FakeResponse()

    monkeypatch.setattr(workbench_api, "urlopen_external", fake_urlopen)
    app = create_app(settings, store=store, runner=WorkbenchRunner(store, settings))
    client = TestClient(app)

    response = client.post(
        "/api/imagegen/generations",
        json={
            "model": "gpt-image-2",
            "prompt": "draw a clean green logo",
            "size": "1024x768",
            "quality": "auto",
            "background": "auto",
            "moderation": "auto",
            "output_format": "png",
            "n": 1,
            "api_base_url": "https://image.example.test",
            "api_key": "request-key",
        },
    )

    assert response.status_code == 200
    assert captured["url"] == "https://image.example.test/v1/images/generations"
    assert captured["authorization"] == "Bearer request-key"
    payload = captured["payload"]
    assert isinstance(payload, dict)
    assert payload["model"] == "gpt-image-2"
    assert payload["size"] == "1024x768"
    assert "resolution" not in payload
    assert "api_base_url" not in payload
    assert "api_key" not in payload


def test_image_generation_endpoint_polls_async_task_result(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.delenv("DRAWAI_IMAGEGEN_API_KEY", raising=False)
    monkeypatch.delenv("OPENAI_API_KEY", raising=False)
    monkeypatch.setattr(workbench_api.time, "sleep", lambda _seconds: None)
    base_config = _base_config(tmp_path)
    store = WorkbenchStore(tmp_path / "workspace")
    settings = _settings(tmp_path, base_config)
    calls: list[dict[str, object]] = []

    class FakeResponse:
        def __init__(self, payload: dict[str, object]) -> None:
            self.payload = payload

        def __enter__(self) -> "FakeResponse":
            return self

        def __exit__(self, exc_type: object, exc: object, tb: object) -> None:
            return None

        def read(self, _: int = -1) -> bytes:
            return json.dumps(self.payload).encode("utf-8")

    def fake_urlopen(request, *, timeout: float):
        calls.append({
            "method": request.get_method(),
            "url": request.full_url,
            "authorization": request.get_header("Authorization"),
            "timeout": timeout,
        })
        if request.get_method() == "POST":
            calls[-1]["payload"] = json.loads(request.data.decode("utf-8"))
            return FakeResponse({"code": 200, "data": [{"status": "submitted", "task_id": "task_abc123"}]})
        return FakeResponse({
            "code": 200,
            "data": {
                "id": "task_abc123",
                "status": "completed",
                "progress": 100,
                "result": {
                    "images": [
                        {
                            "url": ["https://cdn.example.test/final.png"],
                            "expires_at": 1776928569,
                        }
                    ]
                },
            },
        })

    monkeypatch.setattr(workbench_api, "urlopen_external", fake_urlopen)
    app = create_app(settings, store=store, runner=WorkbenchRunner(store, settings))
    client = TestClient(app)

    response = client.post(
        "/api/imagegen/generations",
        json={
            "model": "proxy-image-model",
            "prompt": "draw a clean green logo",
            "size": "1024x768",
            "quality": "auto",
            "background": "transparent",
            "moderation": "auto",
            "output_format": "png",
            "stream": True,
            "partial_images": 1,
            "n": 1,
            "api_base_url": "https://image.example.test",
            "api_key": "request-key",
        },
    )

    assert response.status_code == 200
    assert response.json()["data"]["result"]["images"][0]["url"][0] == "https://cdn.example.test/final.png"
    assert [call["method"] for call in calls] == ["POST", "GET"]
    assert calls[0]["url"] == "https://image.example.test/v1/images/generations"
    assert calls[1]["url"] == "https://image.example.test/v1/tasks/task_abc123"
    assert {call["authorization"] for call in calls} == {"Bearer request-key"}
    post_payload = calls[0]["payload"]
    assert isinstance(post_payload, dict)
    assert post_payload["model"] == "proxy-image-model"
    assert post_payload["size"] == "1024x768"
    assert post_payload["background"] == "transparent"
    assert post_payload["stream"] is True
    assert post_payload["partial_images"] == 1


def test_image_generation_endpoint_can_use_codex_provider(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.delenv("DRAWAI_IMAGEGEN_API_KEY", raising=False)
    monkeypatch.delenv("OPENAI_API_KEY", raising=False)
    base_config = _base_config(tmp_path)
    store = WorkbenchStore(tmp_path / "workspace")
    settings = _settings(tmp_path, base_config)
    captured: list[dict[str, object]] = []

    def fake_invoke_codex_python_sdk_imagegen(**kwargs):
        captured.append(kwargs)
        output_dir = Path(kwargs["output_dir"])
        output_stem = str(kwargs["output_stem"])
        output_dir.mkdir(parents=True, exist_ok=True)
        image_path = output_dir / f"{output_stem}.png"
        Image.new("RGB", (12, 8), (20, 180, 90)).save(image_path)
        image_bytes = image_path.read_bytes()
        image = CodexGeneratedImage(
            image_id=output_stem,
            status="completed",
            path=image_path,
            source_path=str(image_path),
            revised_prompt="codex revised prompt",
            mime_type="image/png",
            width=12,
            height=8,
            bytes=len(image_bytes),
            sha256="test-sha",
        )
        return CodexImageGenResult(
            schema="drawai.codex_python_sdk_imagegen_result.v1",
            runner="codex_python_sdk_imagegen",
            task_name=str(kwargs["task_name"]),
            prompt=str(kwargs["prompt"]),
            final_response='{"generated": true}',
            output_dir=output_dir,
            trace_path=None,
            archive_dir=output_dir / "codex_session_log",
            images=(image,),
            operation="generate",
        )

    monkeypatch.setattr(workbench_api, "invoke_codex_python_sdk_imagegen", fake_invoke_codex_python_sdk_imagegen)
    app = create_app(settings, store=store, runner=WorkbenchRunner(store, settings))
    client = TestClient(app)

    response = client.post(
        "/api/imagegen/generations",
        json={
            "provider": "codex",
            "model": "gpt-image-2",
            "prompt": "draw a clean green logo",
            "size": "2048x1152",
            "quality": "high",
            "background": "transparent",
            "moderation": "auto",
            "output_format": "png",
            "n": 2,
        },
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["provider"] == "codex"
    assert len(payload["data"]) == 2
    assert payload["data"][0]["b64_json"]
    assert payload["data"][0]["size"] == "12x8"
    assert len(captured) == 2
    assert all(call["runtime_config"]["timeout_seconds"] == 300.0 for call in captured)
    assert all("model_name" not in call["runtime_config"] for call in captured)
    first_prompt = str(captured[0]["prompt"])
    assert "Primary request: draw a clean green logo" in first_prompt
    assert "Requested size/aspect: 2048x1152" in first_prompt
    assert "Quality preference: high" in first_prompt
    assert "Background preference: transparent" in first_prompt
    assert "Requested image count: 2" in first_prompt
    assert "image 1 of 2" in first_prompt
    assert "image 2 of 2" in str(captured[1]["prompt"])


def test_image_edit_endpoint_uses_codex_local_source(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    base_config = _base_config(tmp_path)
    store = WorkbenchStore(tmp_path / "workspace")
    settings = _settings(tmp_path, base_config)
    source_image = tmp_path / "source.png"
    Image.new("RGB", (10, 10), (120, 40, 200)).save(source_image)
    captured: dict[str, object] = {}

    def fake_invoke_codex_python_sdk_image_edit(**kwargs):
        captured.update(kwargs)
        output_dir = Path(kwargs["output_dir"])
        output_dir.mkdir(parents=True, exist_ok=True)
        image_path = output_dir / "codex-edited.png"
        Image.new("RGB", (10, 10), (0, 255, 255)).save(image_path)
        image_bytes = image_path.read_bytes()
        image = CodexGeneratedImage(
            image_id="ig_edit",
            status="completed",
            path=image_path,
            source_path=str(image_path),
            revised_prompt="edit revised prompt",
            mime_type="image/png",
            width=10,
            height=10,
            bytes=len(image_bytes),
            sha256="test-sha",
        )
        return CodexImageGenResult(
            schema="drawai.codex_python_sdk_imagegen_result.v1",
            runner="codex_python_sdk_imagegen",
            task_name=str(kwargs["task_name"]),
            prompt=str(kwargs["prompt"]),
            final_response='{"edited": true}',
            output_dir=output_dir,
            trace_path=None,
            archive_dir=output_dir / "codex_session_log",
            images=(image,),
            operation="edit",
            source_image_path=Path(kwargs["source_image_path"]).resolve(),
        )

    monkeypatch.setattr(workbench_api, "invoke_codex_python_sdk_image_edit", fake_invoke_codex_python_sdk_image_edit)
    app = create_app(settings, store=store, runner=WorkbenchRunner(store, settings))
    client = TestClient(app)

    response = client.post(
        "/api/imagegen/edits",
        json={
            "provider": "codex",
            "source_image_path": str(source_image),
            "prompt": "change the circle to cyan",
            "quality": "high",
            "background": "auto",
        },
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["provider"] == "codex"
    assert payload["data"][0]["operation"] == "edit"
    assert payload["data"][0]["b64_json"]
    assert Path(captured["source_image_path"]) == source_image
    edit_prompt = str(captured["prompt"])
    assert "Primary edit request: change the circle to cyan" in edit_prompt
    assert "Quality preference: high" in edit_prompt
    assert "Edit only the supplied image input" in edit_prompt


def test_api_accepts_single_local_image_path(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    source = tmp_path / "single.png"
    Image.new("RGB", (24, 24), "white").save(source)
    settings = _settings(tmp_path, base_config)
    runner = WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor)
    app = create_app(settings, store=store, runner=runner)
    client = TestClient(app)

    response = client.post(
        "/api/batches",
        json={
            "name": "single image batch",
            "input_mode": "local_dir",
            "local_dir": str(source),
            "auto_run_svg_after_analysis": False,
            "max_concurrent_cases": 1,
            "base_config_path": str(base_config),
        },
    )

    assert response.status_code == 200
    assert response.json()["cases"][0]["name"] == "single.png"
    runner.wait_for_idle(timeout=5)
    batch_id = response.json()["batch"]["batch_id"]
    assert client.get(f"/api/batches/{batch_id}").json()["batch"]["status"] == "waiting_review"


def test_api_marks_batch_failed_when_local_path_has_no_images(tmp_path: Path) -> None:
    store = WorkbenchStore(tmp_path / "workspace")
    base_config = _base_config(tmp_path)
    bad_source = tmp_path / "notes.txt"
    bad_source.write_text("not an image", encoding="utf-8")
    settings = _settings(tmp_path, base_config)
    runner = WorkbenchRunner(store, settings, stage_executor=_deterministic_stage_executor)
    app = create_app(settings, store=store, runner=runner)
    client = TestClient(app)

    response = client.post(
        "/api/batches",
        json={
            "name": "bad batch",
            "input_mode": "local_dir",
            "local_dir": str(bad_source),
            "auto_run_svg_after_analysis": False,
            "max_concurrent_cases": 1,
            "base_config_path": str(base_config),
        },
    )

    assert response.status_code == 400
    batches = client.get("/api/batches").json()["batches"]
    assert batches[0]["status"] == "failed"
    assert "supported image" in batches[0]["error_message"]


def _settings(tmp_path: Path, base_config: Path) -> WorkbenchSettings:
    return WorkbenchSettings(
        workspace=tmp_path / "workspace",
        default_config=base_config,
        max_concurrent_cases=2,
    )


def _wait_for_resource_activity(
    runner: WorkbenchRunner,
    resource: str,
    *,
    queued: int,
    running: int,
    timeout: float,
) -> bool:
    deadline = time.monotonic() + timeout
    while time.monotonic() < deadline:
        activity = runner.resource_activity()[resource]
        if activity["queued"] == queued and activity["running"] == running:
            return True
        time.sleep(0.01)
    return False


def _base_config(tmp_path: Path) -> Path:
    path = tmp_path / "base.yaml"
    path.write_text(
        "\n".join(
            [
                "input:",
                f"  image: {tmp_path / 'missing.png'}",
                f"  output_dir: {tmp_path / 'out'}",
                "svg_to_ppt:",
                "  enabled: true",
                "  export_pptx: true",
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    return path


def _deterministic_stage_executor(case, stage: str) -> None:
    root = Path(case.run_root)
    if stage == "prepare":
        (root / "inputs").mkdir(parents=True, exist_ok=True)
        Image.new("RGB", (24, 24), "white").save(root / "inputs" / "figure.png")
        Image.new("RGB", (24, 24), "white").save(root / "inputs" / "original.png")
        _write_json(root / "inputs" / "source_metadata.json", {"width": 24, "height": 24})
    elif stage == "parse_elements":
        _write_json(root / "reports" / "parser_outputs" / "element_candidates.json", {"candidates": []})
    elif stage == "fuse_elements":
        _write_json(root / "trace" / "v2_fusion_trace.json", {"schema": "drawai.v2.fusion_trace.v1"})
    elif stage in {"refine_elements", "asset_analyze"}:
        _write_json(root / "trace" / "v2_refine_trace.json", {"schema": "drawai.v2.refine_trace.v1"})
        _write_json(
            root / "reports" / "element_analysis_codex" / "element_analysis.json",
            {
                "schema": "drawai.codex_element_analysis.v1",
                "elements": [
                    {
                        "box_id": "B001",
                        "source_candidate_ids": ["B001"],
                        "bbox": [1, 1, 10, 10],
                        "category": "crop",
                        "type": "image",
                    },
                    {
                        "box_id": "B002",
                        "source_candidate_ids": ["B002"],
                        "bbox": [11, 1, 20, 10],
                        "category": "svg_self_draw",
                        "type": "arrow",
                    },
                ],
            },
        )
    elif stage == "plan_assets":
        _write_minimal_v2_package(root, case.case_id)
    elif stage == "process_assets":
        _write_minimal_v2_package(root, case.case_id)
    elif stage in {"compose_svg", "svg"}:
        (root / "svg").mkdir(parents=True, exist_ok=True)
        (root / "svg" / "semantic.svg").write_text(
            '<svg xmlns="http://www.w3.org/2000/svg" width="24" height="24" viewBox="0 0 24 24"></svg>\n',
            encoding="utf-8",
        )
        Image.new("RGB", (24, 24), "white").save(root / "svg" / "rendered.png")
        _write_json(root / "reports" / "svg_validation_report.json", {"status": "ok"})
    elif stage == "export":
        (root / "svg_to_ppt").mkdir(parents=True, exist_ok=True)
        (root / "svg_to_ppt" / "semantic.svg_to_ppt.pptx").write_bytes(b"pptx")
        _write_json(
            root / "reports" / "svg_to_ppt_export_report.json",
            {
                "status": "ok",
                "export_backend": "drawai_native_shapes",
                "editable_surface": "native_shapes",
                "requested_export_mode": "native_shapes",
                "effective_export_mode": "native_shapes",
                "export_mode": "native_shapes",
            },
        )


def _write_minimal_v2_package(root: Path, case_id: str) -> None:
    root.mkdir(parents=True, exist_ok=True)
    figure_path = root / "inputs" / "figure.png"
    figure_path.parent.mkdir(parents=True, exist_ok=True)
    if not figure_path.exists():
        Image.new("RGB", (24, 24), "white").save(figure_path)
    _write_json(root / "inputs" / "source_metadata.json", {"width": 24, "height": 24})
    plan = ElementPlan(
        element_id="E001",
        source_candidate_ids=("fixture:E001",),
        element_type="picture",
        bbox=(2.0, 2.0, 12.0, 12.0),
        geometry={"kind": "bbox", "bbox": [2, 2, 12, 12]},
        z_order=1,
        confidence="high",
        processing_intent=ProcessingIntent(
            object_type="picture",
            processing_type="crop",
            parameters={},
        ),
        review_status="deterministic",
        created_by_stage="plan_assets",
        change_reason="Workbench test fixture.",
    )
    asset_package = AssetPackage.empty(
        asset_id="A001",
        element_id="E001",
        processor_type="crop",
    )
    write_element_plan(root, plan)
    write_asset_package(root, asset_package)
    _write_json(
        root / "drawai_package.json",
        {
            "schema": RUN_PACKAGE_SCHEMA,
            "run_id": case_id,
            "root": str(root),
            "source_image": str(figure_path),
            "canvas": {"width": 24, "height": 24},
            "created_at": "2026-06-18T00:00:00Z",
            "metadata": {"last_stage": "plan_assets", "v2_enabled": True},
            "elements": [plan.to_dict()],
            "asset_packages": [asset_package.to_dict()],
        },
    )


def _failing_stage_executor(case, stage: str) -> None:
    if stage == "parse_elements":
        raise RuntimeError("detector unavailable")
    _deterministic_stage_executor(case, stage)


def _export_failing_stage_executor(case, stage: str) -> None:
    if stage == "export":
        raise RuntimeError("export unavailable")
    _deterministic_stage_executor(case, stage)


def _write_single_text_slide_pptx(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_layout)
    box = slide.shapes.add_textbox(914400, 914400, 5486400, 914400)
    box.text = text
    presentation.save(path)


def _write_json(path: Path, payload: object) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


class FakeApiRmbgClient:
    def remove_background(self, image: Image.Image, output_name: str, **_: object) -> RmbgResult:
        result = image.convert("RGBA")
        for x in range(result.width):
            for y in range(result.height):
                r, g, b, a = result.getpixel((x, y))
                if r > 240 and g > 240 and b > 240:
                    result.putpixel((x, y), (255, 255, 255, 0))
                else:
                    result.putpixel((x, y), (r, g, b, a))
        return RmbgResult(image=result, artifacts={"runtime": "fake_api"}, elapsed_ms=5.0)
