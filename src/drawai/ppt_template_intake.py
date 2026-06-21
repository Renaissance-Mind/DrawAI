from __future__ import annotations

from copy import deepcopy
from dataclasses import dataclass
import json
from pathlib import Path
import re
from typing import Any, Mapping, Sequence
import zipfile
import xml.etree.ElementTree as ET


PPT_TEMPLATE_SPEC_SCHEMA = "drawai.ppt_template_spec.v1"
REFERENCE_STYLE_SPEC_SCHEMA = "drawai.reference_style_spec.v1"
PPT_TEMPLATE_FILL_PLAN_SCHEMA = "drawai.ppt_template_fill_plan.v1"

EMU_PER_INCH = 914400
DEFAULT_PPTX_INTAKE_LIMITATIONS = [
    "python-pptx cannot render slide thumbnails by itself; thumbnail_path is left empty unless another renderer is wired.",
    "Slide Master inheritance is summarized from each slide's bound layout, not fully expanded from all master XML.",
    "Native chart/table styling is inventoried coarsely in this first-phase intake.",
]


class PptTemplateIntakeError(ValueError):
    """Raised when a PPTX template cannot be parsed into a DrawAI template spec."""


def intake_ppt_template(
    pptx_path: str | Path,
    *,
    thumbnail_dir: str | Path | None = None,
) -> dict[str, Any]:
    """Parse a PPTX template/example deck into a first-phase DrawAI template spec.

    The output is intentionally close to PPT-master's slide-library contract:
    it inventories canvas, theme identity, per-slide replaceable slots, native
    tables/charts/pictures, and enough geometry for a later template-fill plan.
    """

    source = Path(pptx_path).expanduser().resolve(strict=False)
    if not source.is_file():
        raise PptTemplateIntakeError(f"PPTX template does not exist: {source}")

    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover - import failure depends on local env
        raise PptTemplateIntakeError("python-pptx is required for PPTX template intake") from exc

    prs = Presentation(str(source))
    width_emu = int(prs.slide_width)
    height_emu = int(prs.slide_height)
    theme = _extract_theme_from_pptx(source)

    layouts: list[dict[str, Any]] = []
    observed_colors: set[str] = set()
    observed_fonts: set[str] = set()

    for slide_index, slide in enumerate(prs.slides):
        layout_name = _safe_layout_name(slide)
        slots: list[dict[str, Any]] = []
        tables: list[dict[str, Any]] = []
        charts: list[dict[str, Any]] = []
        pictures: list[dict[str, Any]] = []

        for shape in slide.shapes:
            geometry = _shape_geometry(shape, width_emu=width_emu, height_emu=height_emu)
            colors = _shape_colors(shape)
            observed_colors.update(colors.values())
            font_names = _shape_font_names(shape)
            observed_fonts.update(font_names)

            if _shape_has_table(shape):
                tables.append(_table_inventory(shape, geometry=geometry, slide_index=slide_index))
                continue
            if _shape_has_chart(shape):
                charts.append(_chart_inventory(shape, geometry=geometry, slide_index=slide_index))
                continue
            if _shape_is_picture(shape):
                pictures.append(_picture_inventory(shape, geometry=geometry, slide_index=slide_index))
                continue
            if _shape_has_text(shape):
                slots.append(
                    _text_slot_inventory(
                        shape,
                        geometry=geometry,
                        slide_index=slide_index,
                        slide_size_emu=(width_emu, height_emu),
                    )
                )

        role_guess = _guess_slide_role(slots=slots, tables=tables, charts=charts, pictures=pictures)
        layouts.append(
            {
                "id": f"slide_{slide_index + 1:02d}_{_slugify(layout_name) or 'layout'}",
                "name": layout_name,
                "slide_index": slide_index,
                "role_guess": role_guess,
                "thumbnail_path": _thumbnail_path(thumbnail_dir, source=source, slide_index=slide_index),
                "slots": slots,
                "tables": tables,
                "charts": charts,
                "pictures": pictures,
                "slot_summary": {
                    "text_slot_count": len(slots),
                    "table_count": len(tables),
                    "chart_count": len(charts),
                    "picture_count": len(pictures),
                    "roles": _count_values(slot.get("role") for slot in slots),
                },
            }
        )

    return {
        "schema": PPT_TEMPLATE_SPEC_SCHEMA,
        "source_pptx": str(source),
        "slide_count": len(layouts),
        "slide_size": {
            "width_emu": width_emu,
            "height_emu": height_emu,
            "width_in": round(width_emu / EMU_PER_INCH, 3),
            "height_in": round(height_emu / EMU_PER_INCH, 3),
            "aspect_ratio": round(width_emu / height_emu, 4) if height_emu else None,
        },
        "theme": theme,
        "layouts": layouts,
        "design_tokens": {
            "palette": _sorted_unique(
                [*theme.get("colors", []), *observed_colors],
            ),
            "fonts": _sorted_unique([*theme.get("fonts", []), *observed_fonts]),
            "layout_roles": _count_values(layout.get("role_guess") for layout in layouts),
            "slot_roles": _count_values(
                slot.get("role") for layout in layouts for slot in layout.get("slots", [])
            ),
        },
        "spec_lock": {
            "lock_canvas": True,
            "lock_slot_geometry": True,
            "lock_native_template": True,
            "lock_theme_identity": bool(theme.get("colors") or theme.get("fonts")),
            "source_pptx": str(source),
            "slot_id_contract": "Use layouts[].slots[].slot_id as the stable replacement target in fill plans.",
        },
        "limitations": list(DEFAULT_PPTX_INTAKE_LIMITATIONS),
    }


def write_ppt_template_spec(
    pptx_path: str | Path,
    output_path: str | Path,
    *,
    thumbnail_dir: str | Path | None = None,
) -> dict[str, Any]:
    spec = intake_ppt_template(pptx_path, thumbnail_dir=thumbnail_dir)
    destination = Path(output_path)
    destination.parent.mkdir(parents=True, exist_ok=True)
    destination.write_text(json.dumps(spec, ensure_ascii=False, indent=2), encoding="utf-8")
    return deepcopy(spec)


def build_slot_schema_preview(template_spec: Mapping[str, Any]) -> dict[str, Any]:
    _assert_schema(template_spec, PPT_TEMPLATE_SPEC_SCHEMA, "template_spec")
    return {
        "schema": "drawai.ppt_template_slot_schema_preview.v1",
        "source_pptx": template_spec.get("source_pptx", ""),
        "slide_size": deepcopy(template_spec.get("slide_size", {})),
        "layouts": [
            {
                "layout_id": layout.get("id", ""),
                "slide_index": layout.get("slide_index"),
                "role_guess": layout.get("role_guess", ""),
                "slot_count": len(layout.get("slots", [])),
                "slots": [
                    {
                        "slot_id": slot.get("slot_id", ""),
                        "shape_id": slot.get("shape_id"),
                        "name": slot.get("name", ""),
                        "role": slot.get("role", ""),
                        "placeholder_type": slot.get("placeholder_type", ""),
                        "capacity_visual_width": slot.get("capacity_visual_width"),
                        "old_text": slot.get("text_summary", ""),
                        "geometry": slot.get("geometry", {}),
                    }
                    for slot in layout.get("slots", [])
                ],
                "tables": deepcopy(layout.get("tables", [])),
                "charts": deepcopy(layout.get("charts", [])),
            }
            for layout in template_spec.get("layouts", [])
        ],
    }


def build_prisma_reference_style_spec(
    *,
    source_image_path: str | Path,
    original_source_image_path: str | Path | None = None,
    template_manifest_path: str | Path | None = None,
) -> dict[str, Any]:
    source = Path(source_image_path).expanduser().resolve(strict=False)
    original = (
        str(Path(original_source_image_path).expanduser().resolve(strict=False))
        if original_source_image_path
        else str(source)
    )
    manifest = (
        str(Path(template_manifest_path).expanduser().resolve(strict=False))
        if template_manifest_path
        else ""
    )
    return {
        "schema": REFERENCE_STYLE_SPEC_SCHEMA,
        "id": "prisma_flow_diagram_reference_spec",
        "name": "PRISMA/Systematic Review Flow Reference Style Spec",
        "source_image_path": str(source),
        "original_source_image_path": original,
        "template_manifest_path": manifest,
        "reference_roles": [
            {
                "role": "layout_reference",
                "weight": 0.45,
                "locked_features": [
                    "two-column flow structure",
                    "left vertical stage-label rail",
                    "top yellow section bars",
                    "white rectangular process nodes",
                    "straight arrows linking nodes",
                ],
            },
            {
                "role": "style_reference",
                "weight": 0.25,
                "locked_features": [
                    "clean academic flowchart tone",
                    "thin black strokes",
                    "high whitespace on white background",
                ],
            },
            {
                "role": "color_reference",
                "weight": 0.15,
                "locked_features": ["yellow header bars", "pale-blue stage labels", "black text and arrows"],
            },
            {
                "role": "typography_reference",
                "weight": 0.1,
                "locked_features": ["compact sans-serif labels", "readable node text", "small but clear stage text"],
            },
            {
                "role": "content_reference",
                "weight": 0.05,
                "locked_features": ["workflow stages and exclusion callouts only"],
                "forbidden_copy": [
                    "do not copy original study counts",
                    "do not copy original medical review labels",
                    "do not preserve source facts unless user supplied them",
                ],
            },
        ],
        "design_tokens": {
            "palette": {
                "background": "#ffffff",
                "header_bar": "#fbbf24",
                "stage_label": "#b7d3ee",
                "node_fill": "#ffffff",
                "node_stroke": "#111827",
                "text": "#111827",
                "arrow": "#111827",
            },
            "layout_archetype": "prisma_systematic_review_flow",
            "geometry_rules": [
                "Place section headers in horizontal yellow bars near the top.",
                "Use left-side vertical labels to mark major stages.",
                "Arrange process boxes in stacked flow columns.",
                "Use side boxes only for exclusions or branch outcomes.",
                "Keep arrows orthogonal or straight; avoid decorative connectors.",
            ],
            "slot_schema": {
                "header_bars": {"type": "text_array", "min_items": 1, "max_items": 3},
                "stage_labels": {"type": "text_array", "min_items": 3, "max_items": 5},
                "flow_nodes": {"type": "text_array", "min_items": 6, "max_items": 14},
                "exclusion_nodes": {"type": "text_array", "min_items": 0, "max_items": 6},
                "final_node": {"type": "text", "required": True},
            },
        },
        "spec_lock": {
            "lock_reference_roles": True,
            "lock_layout_archetype": True,
            "lock_color_family": True,
            "allow_content_replacement": True,
            "allow_text_language_change": True,
            "forbid_source_fact_copy": True,
        },
        "limitations": [
            "This spec is declared from the known PRISMA screenshot and manual inspection; no CV segmentation model is wired yet.",
            "It supports layout/style/color/typography guidance but does not yet produce native PPTX slots from the bitmap alone.",
        ],
    }


def build_template_reference_payload(
    *,
    template_spec: Mapping[str, Any],
    reference_style_spec: Mapping[str, Any],
    user_topic: str,
    language: str = "zh",
) -> dict[str, Any]:
    _assert_schema(template_spec, PPT_TEMPLATE_SPEC_SCHEMA, "template_spec")
    _assert_schema(reference_style_spec, REFERENCE_STYLE_SPEC_SCHEMA, "reference_style_spec")
    user_topic = str(user_topic).strip()
    if not user_topic:
        raise PptTemplateIntakeError("user_topic is required")

    recommended_layout = _choose_layout(template_spec)
    reference_roles = deepcopy(reference_style_spec.get("reference_roles", []))
    prompt = "\n".join(
        [
            f"为主题“{user_topic}”生成一页 PPT 视觉草图，语言={language}。",
            "必须结合 PPT 模板结构和参考图风格，而不是自由发挥单页海报。",
            f"PPT 模板来源: {template_spec.get('source_pptx', '')}",
            f"采用模板页面: {recommended_layout.get('id', '')} / {recommended_layout.get('role_guess', '')}",
            "来自 PPT 模板的硬约束: 画布比例、槽位几何、标题/正文/流程/数据区域的可替换 slot。",
            "来自参考图的硬约束: 黄色顶部栏、左侧阶段标签、白底矩形节点、黑色箭头、流程筛选感。",
            "来自用户主题的内容: 替换为 DrawAI/PPT 图像生成或 AI Agent 工作流验证内容，不复制参考图原文和数字。",
            "输出应保留清晰中文标题、阶段标签和流程节点文字。",
        ]
    )

    return {
        "schema": "drawai.template_reference_generation_payload.v1",
        "operation": "edit",
        "provider": "codex",
        "language": language,
        "user_topic": user_topic,
        "source_image_path": reference_style_spec.get("source_image_path", ""),
        "template_spec_path": template_spec.get("source_pptx", ""),
        "reference_style_spec_id": reference_style_spec.get("id", ""),
        "selected_template_layout": {
            "layout_id": recommended_layout.get("id", ""),
            "slide_index": recommended_layout.get("slide_index"),
            "role_guess": recommended_layout.get("role_guess", ""),
            "slot_ids": [slot.get("slot_id", "") for slot in recommended_layout.get("slots", [])],
        },
        "reference_roles": reference_roles,
        "provenance": {
            "from_ppt_template": [
                "slide_size",
                "theme",
                "layouts[].slots[].geometry",
                "layouts[].slots[].role",
                "native table/chart inventory",
            ],
            "from_reference_image": [
                "layout_reference",
                "style_reference",
                "color_reference",
                "typography_reference",
                "content_reference guardrails",
            ],
            "from_user_topic": ["topic", "Chinese replacement text", "workflow semantics"],
        },
        "prompt": prompt,
    }


def create_minimal_fill_plan(template_spec: Mapping[str, Any], *, user_topic: str) -> dict[str, Any]:
    _assert_schema(template_spec, PPT_TEMPLATE_SPEC_SCHEMA, "template_spec")
    topic = str(user_topic).strip() or "DrawAI PPT 图像生成能力验证"
    named_replacements = _named_demo_replacements(topic)
    replacements_by_role = {
        "title": [
            topic,
            "能力分层与交付路径",
            "参考图/模板输入到可编辑重建",
            "验证指标与风险控制",
        ],
        "subtitle": [
            "PPT-master 对齐 demo：保留原生 PPTX 槽位并替换内容",
            "从模板识别、风格锁定到生成后重建",
        ],
        "body": [
            "1. 接收用户 PPT 类型、模板 PPTX、参考图和数据源",
            "2. 抽取 slide_size、theme、slot_schema 与 role_guess",
            "3. 生成 fill_plan，按槽位替换标题、正文、流程与表格",
            "4. 输出可编辑 PPTX，并可继续进入 DrawAI 重建流程",
            "当前阶段已验证 intake + slot preview + 最小填槽。",
            "下一阶段需要完善 Slide Master、图表数据和视觉溢出检查。",
        ],
        "label": ["模板输入", "设计规范", "内容填槽", "质量检查", "可编辑输出"],
        "caption": ["字段来自真实 PPTX 模板与参考图 spec，非纯 prompt。"],
    }
    role_cursors = {role: 0 for role in replacements_by_role}

    slides: list[dict[str, Any]] = []
    for layout in template_spec.get("layouts", []):
        slot_replacements: list[dict[str, Any]] = []
        for slot in layout.get("slots", []):
            role = str(slot.get("role") or "body")
            name = str(slot.get("name") or "").lower()
            text = named_replacements.get(name)
            if text is None:
                pool = replacements_by_role.get(role, replacements_by_role["body"])
                index = role_cursors.get(role, 0)
                role_cursors[role] = index + 1
                text = pool[index % len(pool)]
            slot_replacements.append(
                {
                    "slot_id": slot.get("slot_id", ""),
                    "shape_id": slot.get("shape_id"),
                    "role": role,
                    "old_text": slot.get("text", ""),
                    "new_text": text,
                    "fit_policy": "keep_template_font_and_geometry",
                }
            )
        slides.append(
            {
                "source_slide_index": layout.get("slide_index"),
                "layout_id": layout.get("id", ""),
                "role_guess": layout.get("role_guess", ""),
                "slot_replacements": slot_replacements,
                "table_replacements": [],
                "chart_replacements": [],
            }
        )

    return {
        "schema": PPT_TEMPLATE_FILL_PLAN_SCHEMA,
        "source_pptx": template_spec.get("source_pptx", ""),
        "user_topic": topic,
        "slides": slides,
        "check_policy": {
            "capacity_basis": "capacity_visual_width + role_guess",
            "font_policy": "do_not_shrink_by_default",
            "overflow_policy": "rewrite_or_choose_another_layout",
        },
        "limitations": [
            "This demo fill plan replaces existing text slots in slide order; it does not clone/reorder pages yet.",
            "Table/chart data replacement is inventoried but not applied in this first fill demo.",
        ],
    }


def _named_demo_replacements(topic: str) -> dict[str, str]:
    """Readable replacements for the demo deck's intentionally named slots.

    Real user templates will rely on role/capacity heuristics and a generated fill
    plan; this map keeps the built-in showcase PPTX from looking like a random
    role-fill smoke test.
    """

    return {
        "slot_header_bar": "PPTX 模板输入    design_spec / spec_lock",
        "slot_cover_title": topic,
        "slot_cover_subtitle": "真实 PPTX 模板被解析为可复用 slot_schema，再按 fill_plan 写回原生可编辑元素。",
        "slot_metric_1": "Template\nPPTX",
        "slot_metric_2": "Slots\n23",
        "slot_metric_3": "Output\nEditable",
        "slot_cover_note": "不是整页图片，而是原生 PPTX",
        "slot_content_title": "PPT-master 对齐的数据流",
        "slot_card_1": "1. Template Deck\n接收用户提供的 PPTX / 示例 deck，保留画布、主题和布局槽位。",
        "slot_card_2": "2. Template Spec\n抽取 slide_size、theme、slots、tables、charts 与 role_guess。",
        "slot_card_3": "3. Fill Plan\n用稳定 slot_id 替换内容，并为后续 fit-check 预留容量依据。",
        "slot_content_caption": "这一步已经不是 prompt-only：模板结构来自真实 PPTX，参考图风格来自 typed style spec。",
        "slot_flow_title": "参考图 / PPT 模板 / 数据源的合并路径",
        "slot_flow_step_1": "模板 PPTX",
        "slot_flow_step_2": "Style Spec",
        "slot_flow_step_3": "Fill Plan",
        "slot_flow_step_4": "Editable PPTX",
        "slot_flow_note": "spec_lock 控制颜色、字体、槽位几何；fill_plan 只替换允许变动的内容。",
        "slot_data_title": "第一阶段可验证能力",
        "slot_data_caption": "表格和图表现在可被 intake 识别；数据写回将在下一阶段接入。",
    }


def apply_minimal_fill_plan(
    source_pptx_path: str | Path,
    fill_plan: Mapping[str, Any],
    output_pptx_path: str | Path,
) -> dict[str, Any]:
    _assert_schema(fill_plan, PPT_TEMPLATE_FILL_PLAN_SCHEMA, "fill_plan")
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover - import failure depends on local env
        raise PptTemplateIntakeError("python-pptx is required for PPTX template fill") from exc

    source = Path(source_pptx_path).expanduser().resolve(strict=False)
    if not source.is_file():
        raise PptTemplateIntakeError(f"PPTX template does not exist: {source}")

    prs = Presentation(str(source))
    replaced = 0
    missing: list[dict[str, Any]] = []
    for plan_slide in fill_plan.get("slides", []):
        slide_index = plan_slide.get("source_slide_index")
        if not isinstance(slide_index, int) or slide_index < 0 or slide_index >= len(prs.slides):
            missing.append({"source_slide_index": slide_index, "reason": "slide index out of range"})
            continue
        slide = prs.slides[slide_index]
        shape_by_id = {getattr(shape, "shape_id", None): shape for shape in slide.shapes}
        for replacement in plan_slide.get("slot_replacements", []):
            shape_id = replacement.get("shape_id")
            shape = shape_by_id.get(shape_id)
            if shape is None or not _shape_has_text(shape):
                missing.append(
                    {
                        "source_slide_index": slide_index,
                        "slot_id": replacement.get("slot_id", ""),
                        "shape_id": shape_id,
                        "reason": "shape not found or not text",
                    }
                )
                continue
            shape.text = str(replacement.get("new_text", ""))
            replaced += 1

    output = Path(output_pptx_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return {
        "schema": "drawai.ppt_template_fill_result.v1",
        "source_pptx": str(source),
        "output_pptx": str(output.resolve(strict=False)),
        "replaced_slot_count": replaced,
        "missing_targets": missing,
    }


def save_json(payload: Mapping[str, Any], path: str | Path) -> None:
    destination = Path(path)
    destination.parent.mkdir(parents=True, exist_ok=True)
    destination.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def _assert_schema(payload: Mapping[str, Any], schema: str, label: str) -> None:
    if payload.get("schema") != schema:
        raise PptTemplateIntakeError(f"{label} must use schema {schema}")


def _extract_theme_from_pptx(pptx_path: Path) -> dict[str, Any]:
    theme = {"colors": [], "fonts": [], "raw_color_scheme": {}, "raw_font_scheme": {}}
    try:
        with zipfile.ZipFile(pptx_path) as archive:
            theme_names = [name for name in archive.namelist() if name.startswith("ppt/theme/") and name.endswith(".xml")]
            if not theme_names:
                return theme
            xml = archive.read(theme_names[0])
    except Exception:
        return theme

    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    try:
        root = ET.fromstring(xml)
    except ET.ParseError:
        return theme

    color_scheme: dict[str, str] = {}
    clr_scheme = root.find(".//a:clrScheme", ns)
    if clr_scheme is not None:
        for child in list(clr_scheme):
            color_name = _local_name(child.tag)
            color_value = ""
            srgb = child.find(".//a:srgbClr", ns)
            sys_clr = child.find(".//a:sysClr", ns)
            if srgb is not None:
                color_value = "#" + str(srgb.attrib.get("val", "")).upper()
            elif sys_clr is not None:
                color_value = "#" + str(sys_clr.attrib.get("lastClr", "")).upper()
            if color_value and len(color_value) == 7:
                color_scheme[color_name] = color_value

    font_scheme: dict[str, str] = {}
    for parent_name, xpath in {
        "major_latin": ".//a:majorFont/a:latin",
        "minor_latin": ".//a:minorFont/a:latin",
        "major_ea": ".//a:majorFont/a:ea",
        "minor_ea": ".//a:minorFont/a:ea",
    }.items():
        node = root.find(xpath, ns)
        if node is not None and node.attrib.get("typeface"):
            font_scheme[parent_name] = node.attrib["typeface"]

    theme["raw_color_scheme"] = color_scheme
    theme["raw_font_scheme"] = font_scheme
    theme["colors"] = _sorted_unique(color_scheme.values())
    theme["fonts"] = _sorted_unique(value for value in font_scheme.values() if value)
    return theme


def _safe_layout_name(slide: Any) -> str:
    try:
        return str(slide.slide_layout.name or "layout")
    except Exception:
        return "layout"


def _shape_geometry(shape: Any, *, width_emu: int, height_emu: int) -> dict[str, Any]:
    left = int(getattr(shape, "left", 0) or 0)
    top = int(getattr(shape, "top", 0) or 0)
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    return {
        "left_emu": left,
        "top_emu": top,
        "width_emu": width,
        "height_emu": height,
        "left_pct": round(left / width_emu, 4) if width_emu else None,
        "top_pct": round(top / height_emu, 4) if height_emu else None,
        "width_pct": round(width / width_emu, 4) if width_emu else None,
        "height_pct": round(height / height_emu, 4) if height_emu else None,
        "left_in": round(left / EMU_PER_INCH, 3),
        "top_in": round(top / EMU_PER_INCH, 3),
        "width_in": round(width / EMU_PER_INCH, 3),
        "height_in": round(height / EMU_PER_INCH, 3),
    }


def _shape_has_text(shape: Any) -> bool:
    return bool(getattr(shape, "has_text_frame", False))


def _shape_has_table(shape: Any) -> bool:
    return bool(getattr(shape, "has_table", False))


def _shape_has_chart(shape: Any) -> bool:
    return bool(getattr(shape, "has_chart", False))


def _shape_is_picture(shape: Any) -> bool:
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    except Exception:
        return False


def _text_slot_inventory(
    shape: Any,
    *,
    geometry: Mapping[str, Any],
    slide_index: int,
    slide_size_emu: tuple[int, int],
) -> dict[str, Any]:
    shape_id = int(getattr(shape, "shape_id", 0) or 0)
    text = _shape_text(shape)
    font_size_pt = _shape_font_size_pt(shape)
    font_size_px = round(font_size_pt * 96 / 72, 1) if font_size_pt else None
    role = _guess_text_role(shape, geometry=geometry, text=text, font_size_pt=font_size_pt)
    paragraph_count = len(getattr(getattr(shape, "text_frame", None), "paragraphs", []) or [])
    return {
        "slot_id": f"s{slide_index + 1:02d}_shape_{shape_id}",
        "shape_id": shape_id,
        "name": str(getattr(shape, "name", "") or ""),
        "kind": "text",
        "role": role,
        "placeholder_type": _placeholder_type(shape),
        "geometry": dict(geometry),
        "paragraph_count": paragraph_count,
        "text": text,
        "text_summary": _summarize_text(text),
        "text_metrics": {
            "font_size_pt": font_size_pt,
            "font_size_px": font_size_px,
            "font_names": _sorted_unique(_shape_font_names(shape)),
        },
        "capacity_visual_width": _estimate_capacity(geometry, font_size_px=font_size_px),
        "style": _shape_style(shape),
    }


def _table_inventory(shape: Any, *, geometry: Mapping[str, Any], slide_index: int) -> dict[str, Any]:
    table = shape.table
    cells: list[dict[str, Any]] = []
    for row_index, row in enumerate(table.rows):
        for col_index, cell in enumerate(row.cells):
            cells.append({"row": row_index, "col": col_index, "text": str(cell.text or "")})
    shape_id = int(getattr(shape, "shape_id", 0) or 0)
    return {
        "table_id": f"s{slide_index + 1:02d}_table_{shape_id}",
        "shape_id": shape_id,
        "name": str(getattr(shape, "name", "") or ""),
        "row_count": len(table.rows),
        "column_count": len(table.columns),
        "geometry": dict(geometry),
        "cells": cells,
    }


def _chart_inventory(shape: Any, *, geometry: Mapping[str, Any], slide_index: int) -> dict[str, Any]:
    shape_id = int(getattr(shape, "shape_id", 0) or 0)
    chart = shape.chart
    chart_type = ""
    try:
        chart_type = str(chart.chart_type)
    except Exception:
        pass
    series = []
    try:
        series = [str(item.name) for item in chart.series]
    except Exception:
        pass
    return {
        "chart_id": f"s{slide_index + 1:02d}_chart_{shape_id}",
        "shape_id": shape_id,
        "name": str(getattr(shape, "name", "") or ""),
        "chart_type": chart_type,
        "series_names": series,
        "geometry": dict(geometry),
    }


def _picture_inventory(shape: Any, *, geometry: Mapping[str, Any], slide_index: int) -> dict[str, Any]:
    shape_id = int(getattr(shape, "shape_id", 0) or 0)
    return {
        "picture_id": f"s{slide_index + 1:02d}_picture_{shape_id}",
        "shape_id": shape_id,
        "name": str(getattr(shape, "name", "") or ""),
        "geometry": dict(geometry),
    }


def _shape_text(shape: Any) -> str:
    try:
        return "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
    except Exception:
        try:
            return str(shape.text or "").strip()
        except Exception:
            return ""


def _shape_font_size_pt(shape: Any) -> float | None:
    try:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                size = run.font.size
                if size is not None:
                    value = getattr(size, "pt", None)
                    if value:
                        return round(float(value), 1)
    except Exception:
        return None
    return None


def _shape_font_names(shape: Any) -> list[str]:
    fonts: set[str] = set()
    try:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                name = run.font.name
                if name:
                    fonts.add(str(name))
    except Exception:
        pass
    return sorted(fonts)


def _shape_colors(shape: Any) -> dict[str, str]:
    colors: dict[str, str] = {}
    try:
        fill_rgb = getattr(shape.fill.fore_color, "rgb", None)
        if fill_rgb is not None:
            colors["fill"] = "#" + str(fill_rgb).upper()
    except Exception:
        pass
    try:
        line_rgb = getattr(shape.line.color, "rgb", None)
        if line_rgb is not None:
            colors["line"] = "#" + str(line_rgb).upper()
    except Exception:
        pass
    return colors


def _shape_style(shape: Any) -> dict[str, Any]:
    return {"colors": _shape_colors(shape), "font_names": _shape_font_names(shape)}


def _placeholder_type(shape: Any) -> str:
    try:
        if not shape.is_placeholder:
            return ""
        return str(shape.placeholder_format.type)
    except Exception:
        return ""


def _guess_text_role(
    shape: Any,
    *,
    geometry: Mapping[str, Any],
    text: str,
    font_size_pt: float | None,
) -> str:
    name = str(getattr(shape, "name", "") or "").lower()
    placeholder = _placeholder_type(shape).lower()
    text_len = len(text)
    top = float(geometry.get("top_pct") or 0)
    height = float(geometry.get("height_pct") or 0)
    width = float(geometry.get("width_pct") or 0)
    if "title" in name or "title" in placeholder or (top < 0.22 and (font_size_pt or 0) >= 24):
        return "title"
    if "subtitle" in name or "subtitle" in placeholder:
        return "subtitle"
    if "caption" in name or height < 0.08 and text_len <= 60:
        return "caption"
    if "label" in name or (width < 0.22 and height < 0.18 and text_len <= 30):
        return "label"
    if text_len <= 24 and height <= 0.16:
        return "label"
    return "body"


def _guess_slide_role(
    *,
    slots: Sequence[Mapping[str, Any]],
    tables: Sequence[Mapping[str, Any]],
    charts: Sequence[Mapping[str, Any]],
    pictures: Sequence[Mapping[str, Any]],
) -> str:
    title_count = sum(1 for slot in slots if slot.get("role") == "title")
    body_count = sum(1 for slot in slots if slot.get("role") == "body")
    label_count = sum(1 for slot in slots if slot.get("role") == "label")
    if charts or tables:
        return "data_page"
    if title_count and body_count == 0 and label_count <= 3:
        return "cover"
    if label_count >= 4 and body_count <= 4:
        return "process_or_timeline"
    if pictures and body_count <= 3:
        return "visual_feature"
    if body_count >= 4:
        return "multi_card_content"
    return "content"


def _estimate_capacity(geometry: Mapping[str, Any], *, font_size_px: float | None) -> int:
    width_px = float(geometry.get("width_in") or 0) * 96
    height_px = float(geometry.get("height_in") or 0) * 96
    font_px = font_size_px or 18.0
    if width_px <= 0 or height_px <= 0:
        return 0
    line_count = max(1, int(height_px / (font_px * 1.35)))
    chars_per_line = max(4, int(width_px / (font_px * 0.58)))
    return int(line_count * chars_per_line)


def _thumbnail_path(thumbnail_dir: str | Path | None, *, source: Path, slide_index: int) -> str | None:
    if thumbnail_dir is None:
        return None
    path = Path(thumbnail_dir) / f"{source.stem}_slide_{slide_index + 1:02d}.png"
    return str(path)


def _choose_layout(template_spec: Mapping[str, Any]) -> dict[str, Any]:
    layouts = [dict(layout) for layout in template_spec.get("layouts", [])]
    if not layouts:
        return {}
    for preferred_role in ("process_or_timeline", "multi_card_content", "content", "cover"):
        for layout in layouts:
            if layout.get("role_guess") == preferred_role:
                return layout
    return layouts[0]


def _summarize_text(text: str, *, limit: int = 90) -> str:
    clean = re.sub(r"\s+", " ", str(text or "")).strip()
    if len(clean) <= limit:
        return clean
    return clean[: limit - 1] + "…"


def _count_values(values: Any) -> dict[str, int]:
    counts: dict[str, int] = {}
    for value in values:
        key = str(value or "")
        if not key:
            continue
        counts[key] = counts.get(key, 0) + 1
    return dict(sorted(counts.items()))


def _sorted_unique(values: Any) -> list[str]:
    return sorted({str(value) for value in values if str(value or "").strip()})


def _slugify(value: str) -> str:
    slug = re.sub(r"[^a-zA-Z0-9]+", "_", value.lower()).strip("_")
    return slug[:48]


def _local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


__all__ = [
    "PPT_TEMPLATE_FILL_PLAN_SCHEMA",
    "PPT_TEMPLATE_SPEC_SCHEMA",
    "REFERENCE_STYLE_SPEC_SCHEMA",
    "PptTemplateIntakeError",
    "apply_minimal_fill_plan",
    "build_prisma_reference_style_spec",
    "build_slot_schema_preview",
    "build_template_reference_payload",
    "create_minimal_fill_plan",
    "intake_ppt_template",
    "save_json",
    "write_ppt_template_spec",
]
